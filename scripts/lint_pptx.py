# -*- coding: utf-8 -*-
"""对 apply_content 后的 PPTX 做规则层抽检（三层混合架构 v2 版）。

对比 template_story.json（v2）与当前 PPTX：
- 占位几何被改动（容差 1 EMU；同页重复 shape_id 跳过）
- content_shape 文本超出 char_limit.max（text_exceeds_char_limit）
- content_shape 段数大幅偏离 paragraph_count（paragraph_count_mismatch，warning）
- 同页多个 content_shape 完全重复文本（duplicate_text_in_slide）
- shape_group 全部成员文本雷同（shape_group_identical）
- is_placeholder_text=true 的 shape 仍保留原模板占位文本（placeholder_not_replaced，warning）
- non_content keep_original 被改写（non_content_modified / decoration_text_lost）
- non_content clear_to_empty 后仍有残留文本（non_content_not_cleared）

CLI:
    python lint_pptx.py <pptx-after> --story template_story.json --out lint_report.json
"""

from __future__ import annotations

import argparse
import json
import sys
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


_DECORATION_ROLES = {"decoration_number", "style_tag", "logo_text",
                     "decoration_micro", "time_marker", "page_meta"}

_CLEAR_TO_EMPTY_ROLES = {"template_sample_text"}


def _walk_shapes(container, path: str = ""):
    iterable = container.shapes if hasattr(container, "shapes") else container
    for shape in iterable:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            sub = f"{path}/{shape.name}" if path else shape.name
            yield from _walk_shapes(shape, path=sub)
        else:
            yield shape, path


def _shape_local_id(shape, path: str) -> str:
    """当前 pptx 中的 shape 本地 id（便于和 story.debug_name 对齐）。"""
    return f"{path}/{shape.name}" if path else shape.name


def _shape_cnvpr_id(shape) -> int | None:
    try:
        return int(shape._element.nvSpPr.cNvPr.get("id"))  # type: ignore[attr-defined]
    except Exception:
        try:
            return int(shape._element.nvGrpSpPr.cNvPr.get("id"))  # type: ignore[attr-defined]
        except Exception:
            try:
                # picture / group fallbacks
                el = shape._element
                cNvPr = el.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr")
                if cNvPr is not None:
                    return int(cNvPr.get("id"))
            except Exception:
                return None
    return None


def _collect_text(shape) -> tuple[str, list[str]]:
    if not getattr(shape, "has_text_frame", False):
        return "", []
    paras = [p.text for p in shape.text_frame.paragraphs]
    return "\n".join(paras), paras


def _count_chars(text: str) -> int:
    return sum(1 for c in (text or "") if not c.isspace())


def _build_story_idx(story: dict) -> dict:
    """key = slide_{i}::sp_{cNvPr.id}（与 apply 层一致），值包含 content/non_content 元信息。"""
    idx: dict[str, dict] = {}
    for s in story.get("slides", []):
        i = s["slide_index"]
        for cs in s.get("content_shapes", []):
            sid = cs["shape_id"]
            idx[sid] = {
                "type": "content",
                "slide": i,
                "bbox": cs.get("bbox"),
                "char_limit": cs.get("char_limit", {"min": 1, "max": 0}),
                "hard_ceiling_chars": cs.get("hard_ceiling_chars"),
                "per_paragraph": cs.get("per_paragraph") or [],
                "paragraph_count": cs.get("paragraph_count", 1),
                "original_text": cs.get("original_text", ""),
                "is_placeholder_text": bool(cs.get("is_placeholder_text")),
                "parent_table_id": cs.get("parent_table_id"),
            }
        for nc in s.get("non_content_shapes", []):
            sid = nc["shape_id"]
            idx[sid] = {
                "type": "non_content",
                "slide": i,
                "role": nc.get("role"),
                "preserve_action": nc.get("preserve_action", "keep_original"),
                "bbox": nc.get("bbox"),
                "original_text": nc.get("current_text", ""),
            }
    return idx


def _shape_groups_by_slide(story: dict) -> dict[int, list[dict]]:
    out: dict[int, list[dict]] = defaultdict(list)
    for s in story.get("slides", []):
        for g in s.get("shape_groups", []):
            out[s["slide_index"]].append(g)
    return out


def _map_live_shape_id(slide_idx: int, shape) -> str:
    cnv = _shape_cnvpr_id(shape)
    if cnv is not None:
        return f"slide_{slide_idx}::sp_{cnv}"
    return f"slide_{slide_idx}::{shape.name}"


def lint(pptx_path: Path, story: dict) -> dict:
    prs = Presentation(str(pptx_path))
    story_idx = _build_story_idx(story)
    groups_by_slide = _shape_groups_by_slide(story)

    findings: list[dict] = []
    stats = {
        "slides": 0,
        "content_shapes": sum(1 for v in story_idx.values() if v["type"] == "content"),
        "non_content_shapes": sum(1 for v in story_idx.values() if v["type"] == "non_content"),
        "filled": 0,
        "cleared": 0,
        "placeholder_not_replaced": 0,
        "text_exceeds_char_limit": 0,
        "paragraph_count_mismatch": 0,
        "duplicate_text_in_slide": 0,
        "shape_group_identical": 0,
        "non_content_modified": 0,
        "non_content_not_cleared": 0,
        "decoration_text_lost": 0,
        "geometry_changed": 0,
        "paragraph_char_exceeds_limit": 0,
        "emphasis_paragraph_overflow": 0,
        "errors": 0,
        "warnings": 0,
    }

    for i, slide in enumerate(prs.slides, start=1):
        stats["slides"] += 1

        # 当前 pptx 里按 shape_id 汇总文本
        live_text_by_sid: dict[str, str] = {}
        live_paras_by_sid: dict[str, list[str]] = {}
        live_shape_by_sid: dict[str, object] = {}

        # 同 id 出现次数（重复 id 跳过几何 diff，以避免 cNvPr.id 冲突带来的误伤）
        sid_counts: dict[str, int] = defaultdict(int)
        for shape, _gp in _walk_shapes(slide):
            sid = _map_live_shape_id(i, shape)
            sid_counts[sid] += 1

        # 表格 cell 的 shape_id 特殊形态: slide_{i}::sp_{tableId}::cell_{r}_{c}
        for shape, _gp in _walk_shapes(slide):
            sid = _map_live_shape_id(i, shape)
            # 表格拆 cell
            if getattr(shape, "has_table", False):
                try:
                    cnv = _shape_cnvpr_id(shape)
                    for r, row in enumerate(shape.table.rows):
                        for c, cell in enumerate(row.cells):
                            cell_sid = f"slide_{i}::sp_{cnv}::cell_{r}_{c}"
                            tf = cell.text_frame
                            paras = [p.text for p in tf.paragraphs]
                            live_paras_by_sid[cell_sid] = paras
                            live_text_by_sid[cell_sid] = "\n".join(paras)
                            live_shape_by_sid[cell_sid] = shape  # 共享 shape 做几何
                except Exception:
                    pass
                continue
            txt, paras = _collect_text(shape)
            live_text_by_sid[sid] = txt
            live_paras_by_sid[sid] = paras
            live_shape_by_sid[sid] = shape

        # 遍历 story_idx 中属于本页的条目
        slide_content_texts: list[dict] = []
        for sid, meta in story_idx.items():
            if meta["slide"] != i:
                continue
            cur_text = live_text_by_sid.get(sid, "")
            cur_paras = live_paras_by_sid.get(sid, [])
            shape = live_shape_by_sid.get(sid)
            cur_stripped = cur_text.strip()

            # 几何 diff（仅对本体 shape 且同页 id 唯一；表格 cell 跳过）
            base_sid = sid.split("::cell_")[0] if "::cell_" in sid else sid
            if (meta.get("bbox") and "::cell_" not in sid
                    and sid_counts.get(base_sid, 0) == 1
                    and shape is not None):
                try:
                    left = int(shape.left or 0)
                    top = int(shape.top or 0)
                    w = int(shape.width or 0)
                    h = int(shape.height or 0)
                except Exception:
                    left = top = w = h = 0
                bbox = meta["bbox"]
                for fld, cur, pre in [
                    ("left_emu", left, int(bbox.get("left_emu", left))),
                    ("top_emu", top, int(bbox.get("top_emu", top))),
                    ("w_emu", w, int(bbox.get("w_emu", w))),
                    ("h_emu", h, int(bbox.get("h_emu", h))),
                ]:
                    if abs(cur - pre) > 1:
                        findings.append({
                            "slide": i, "shape_id": sid,
                            "kind": "geometry_changed",
                            "field": fld, "before": pre, "after": cur,
                        })
                        stats["geometry_changed"] += 1
                        stats["errors"] += 1

            # non_content 检查
            if meta["type"] == "non_content":
                before = (meta.get("original_text") or "").strip()
                action = meta["preserve_action"]
                role = meta.get("role") or ""
                if cur_stripped:
                    stats["filled"] += 1
                else:
                    stats["cleared"] += 1
                if action == "keep_original":
                    if cur_stripped != before:
                        if role in _DECORATION_ROLES:
                            findings.append({
                                "slide": i, "shape_id": sid,
                                "kind": "decoration_text_lost",
                                "role": role,
                                "before": before[:60],
                                "after": cur_stripped[:60],
                            })
                            stats["decoration_text_lost"] += 1
                        else:
                            findings.append({
                                "slide": i, "shape_id": sid,
                                "kind": "non_content_modified",
                                "role": role,
                                "before": before[:60],
                                "after": cur_stripped[:60],
                            })
                            stats["non_content_modified"] += 1
                        stats["errors"] += 1
                elif action == "clear_to_empty":
                    if cur_stripped:
                        findings.append({
                            "slide": i, "shape_id": sid,
                            "kind": "non_content_not_cleared",
                            "role": role,
                            "after": cur_stripped[:60],
                        })
                        stats["non_content_not_cleared"] += 1
                        stats["warnings"] += 1
                continue

            # content_shape 检查
            if cur_stripped:
                stats["filled"] += 1
            char_limit = meta.get("char_limit") or {}
            mc = int(char_limit.get("max", 0))
            if mc and _count_chars(cur_text) > mc:
                findings.append({
                    "slide": i, "shape_id": sid,
                    "kind": "text_exceeds_char_limit",
                    "char_count": _count_chars(cur_text),
                    "max_chars": mc,
                })
                stats["text_exceeds_char_limit"] += 1
                stats["errors"] += 1

            # per_paragraph 逐段 char_limit 检查
            per_paragraph = meta.get("per_paragraph") or []
            if per_paragraph and cur_paras:
                for p_idx, pp in enumerate(per_paragraph):
                    if p_idx >= len(cur_paras):
                        break
                    p_text = cur_paras[p_idx] or ""
                    p_chars = _count_chars(p_text)
                    pl = pp.get("char_limit") or {}
                    p_max = int(pl.get("max", 0))
                    if p_max and p_chars > p_max:
                        kind = ("emphasis_paragraph_overflow"
                                if pp.get("is_emphasis")
                                else "paragraph_char_exceeds_limit")
                        findings.append({
                            "slide": i, "shape_id": sid,
                            "kind": kind,
                            "paragraph_index": p_idx,
                            "char_count": p_chars,
                            "max_chars": p_max,
                            "is_emphasis": bool(pp.get("is_emphasis")),
                            "font_size_pt": pp.get("font_size_pt"),
                        })
                        stats[kind] += 1
                        stats["errors"] += 1

            expected_paras = int(meta.get("paragraph_count") or 1)
            actual_paras = max(1, len([p for p in cur_paras if p is not None]))
            # 只有多段 shape 才严格要求；相差 > 1 视为偏离
            if expected_paras > 1 and abs(actual_paras - expected_paras) > 1:
                findings.append({
                    "slide": i, "shape_id": sid,
                    "kind": "paragraph_count_mismatch",
                    "expected": expected_paras,
                    "actual": actual_paras,
                })
                stats["paragraph_count_mismatch"] += 1
                stats["warnings"] += 1

            # 占位符未替换：is_placeholder_text=true 且文本等于原占位
            if meta.get("is_placeholder_text"):
                before = (meta.get("original_text") or "").strip()
                if cur_stripped and before and cur_stripped == before:
                    findings.append({
                        "slide": i, "shape_id": sid,
                        "kind": "placeholder_not_replaced",
                        "text": before[:60],
                    })
                    stats["placeholder_not_replaced"] += 1
                    stats["warnings"] += 1

            if cur_stripped:
                slide_content_texts.append({
                    "sid": sid, "text": cur_stripped,
                })

        # 页内重复：两个及以上 shape 完全重复（允许在同一 shape_group 内的重复？
        # 按设计，shape_group 本意是"并列但语义不同"，完全一致仍视为错误）
        text_map: dict[str, list[str]] = defaultdict(list)
        for it in slide_content_texts:
            text_map[it["text"]].append(it["sid"])
        for text, sids in text_map.items():
            if len(sids) > 1:
                findings.append({
                    "slide": i,
                    "kind": "duplicate_text_in_slide",
                    "text": text[:60],
                    "shape_ids": sids[:6],
                    "count": len(sids),
                })
                stats["duplicate_text_in_slide"] += 1
                stats["errors"] += 1

        # shape_group 全部雷同
        for g in groups_by_slide.get(i, []):
            members = g.get("member_shape_ids") or []
            texts = [live_text_by_sid.get(sid, "").strip() for sid in members]
            texts = [t for t in texts if t]
            if len(texts) >= 3 and len(set(texts)) == 1:
                findings.append({
                    "slide": i,
                    "kind": "shape_group_identical",
                    "group_id": g.get("group_id"),
                    "shared_text": texts[0][:60],
                    "member_count": len(texts),
                })
                stats["shape_group_identical"] += 1
                stats["errors"] += 1

    return {"ok": stats["errors"] == 0, "stats": stats, "findings": findings}


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("pptx", type=Path)
    ap.add_argument("--story", required=True, type=Path)
    ap.add_argument("--out", required=True, type=Path)
    args = ap.parse_args()

    if not args.pptx.exists():
        print(f"ERROR: pptx not found: {args.pptx}", file=sys.stderr)
        return 1
    story = json.loads(args.story.read_text(encoding="utf-8"))
    if int(story.get("schema_version", 0)) < 2:
        print("ERROR: template_story.json schema_version < 2, 请重新 parse",
              file=sys.stderr)
        return 1
    report = lint(args.pptx, story)
    args.out.parent.mkdir(parents=True, exist_ok=True)
    args.out.write_text(json.dumps(report, ensure_ascii=False, indent=2),
                        encoding="utf-8")
    print(json.dumps({"ok": report["ok"], "out": str(args.out),
                      **report["stats"]}, ensure_ascii=False))
    return 0 if report["ok"] else 2


if __name__ == "__main__":
    sys.exit(main())
