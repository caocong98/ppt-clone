"""把 content_mapping.json (v4) 应用到 PPT。

特性：
- shape_id 用 `slide_{i}::sp_{cNvPr.id}`（rerun 稳定）；table cell 用
  `slide_{i}::sp_{tableId}::cell_{r}_{c}` 单独定位。
- 特殊值：
  "__skip__"  → 完全不动该 shape
  "__clear__" → 清空文字但保留 <a:bodyPr>/<a:lstStyle>/<a:pPr>/<a:rPr>/段落节点
- 多段落支持：value 为 list[str] 时按段落写入并保留每段 <a:pPr>（bullet/缩进/level）；
  原段数 < 新段数时克隆最后一段的 pPr；原段数 > 新段数时清空多余段。
- L5.1 装饰兜底：原文是装饰字符（壹/Ⅰ/A 等）且新值远大于原文 → 强制跳过。
- L5.2 硬字数闸门 `_enforce_char_limit`：
    chars > max * 2          → forced_skipped 拒写
    max < chars <= max * 2   → 截断到 max 并补 "…"
    chars <= max             → 通过
- L5.3 capacity gate：mapping.pending_shrink_requests 非空 → 拒绝落盘
  （--force-apply 仅绕过此闸门，**不**绕过 char_limit 和装饰保护）

CLI:
    python apply_content.py <pptx> --mapping content_mapping.json --out final.pptx \\
        [--story template_story.json] [--force-apply]
"""

from __future__ import annotations

import argparse
import copy
import json
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

try:
    from parse_template_story import _is_ordinal_decoration as _ord_detect  # type: ignore
except Exception:  # pragma: no cover
    _ord_detect = None  # type: ignore


# === 装饰文本保护兜底 ===

def _looks_like_decoration(original_text: str) -> str | None:
    if not original_text:
        return None
    t = original_text.strip()
    if not t or len(t) > 4:
        return None
    if _ord_detect is not None:
        sub = _ord_detect(t)
        if sub is not None:
            return f"ordinal_{sub}"
    if t.isascii() and t.isalpha() and len(t) <= 3:
        return "latin_short"
    return None


def _decoration_protection_triggered(original_text: str, new_value) -> bool:
    kind = _looks_like_decoration(original_text)
    if kind is None:
        return False
    if not isinstance(new_value, str):
        return True
    new_len = len(new_value.strip())
    return new_len >= max(len(original_text.strip()) * 3, 6)


# === 硬字数闸门 ===

ELLIPSIS = "…"


def _count_chars(text: str) -> int:
    if not text:
        return 0
    return sum(1 for c in text if c not in (" ", "\n", "\t", "\r", "\u3000"))


def _enforce_char_limit(value,
                        char_limit: dict | None,
                        *,
                        per_paragraph: list[dict] | None = None,
                        hard_ceiling: int | None = None,
                        ) -> tuple[object, dict | None]:
    """根据 char_limit.max / hard_ceiling_chars + per_paragraph 做硬截断/拒写。

    - str 值：
        chars > hard_ceiling*2 → forced_skipped
        chars > hard_ceiling   → 截断到 hard_ceiling 并补 "…"
        否则通过
    - list 值 + per_paragraph：
        逐段用 per_paragraph[i].char_limit.max 截断；
        is_emphasis 段若 chars > p_max*1.5 → **本段单独保留原文占位**（不波及整 shape）；
        整体 total_chars > sum(hard_ceiling_i)*2 → forced_skipped
    - list 值 + 无 per_paragraph：沿用按比例截断。
    """
    if not char_limit and not hard_ceiling:
        return value, None
    max_chars = int((char_limit or {}).get("max", 0))
    hard = int(hard_ceiling) if hard_ceiling else max_chars
    if hard <= 0:
        hard = max_chars
    effective_truncate = hard if hard > 0 else max_chars
    if effective_truncate <= 0:
        return value, None

    # -------- str --------
    if isinstance(value, str):
        chars = _count_chars(value)
        if chars <= effective_truncate:
            return value, None
        if chars > effective_truncate * 2:
            return value, {
                "action": "forced_skipped",
                "chars": chars,
                "max_chars": max_chars,
                "hard_ceiling": effective_truncate,
                "reason": "chars_exceed_2x_hard_ceiling",
            }
        truncated = _truncate_with_ellipsis(value, effective_truncate)
        return truncated, {
            "action": "truncated",
            "chars_before": chars,
            "chars_after": _count_chars(truncated),
            "max_chars": max_chars,
            "hard_ceiling": effective_truncate,
        }

    # -------- list --------
    if isinstance(value, list):
        total = sum(_count_chars(p) for p in value if isinstance(p, str))
        # per_paragraph 路径（逐段独立处理）
        if per_paragraph:
            new_paragraphs: list[str] = []
            per_para_truncated = 0
            emphasis_protected = 0
            forced_whole_skip = False
            hard_sum = 0
            for i, p in enumerate(value):
                if not isinstance(p, str):
                    new_paragraphs.append(p)
                    continue
                pp = per_paragraph[i] if i < len(per_paragraph) else None
                if not pp:
                    new_paragraphs.append(p)
                    continue
                pl = pp.get("char_limit") or {}
                p_max = int(pl.get("max", 0))
                p_hard = int(pp.get("hard_ceiling_chars") or p_max)
                is_emphasis = bool(pp.get("is_emphasis"))
                p_chars = _count_chars(p)
                hard_sum += p_hard

                if p_hard <= 0:
                    new_paragraphs.append(p)
                    continue

                # emphasis 段：严重超限 → 本段"放弃 LLM 文本，保留原样"占位标记
                if is_emphasis and p_max > 0 and p_chars > int(p_max * 1.5):
                    # 用 __skip_paragraph__ 标记；下游 _apply_paragraphs 识别后跳过该段
                    new_paragraphs.append("__skip_paragraph__")
                    emphasis_protected += 1
                    continue

                if p_chars > p_hard:
                    truncated = _truncate_with_ellipsis(p, p_hard)
                    new_paragraphs.append(truncated)
                    per_para_truncated += 1
                else:
                    new_paragraphs.append(p)

            # 整 shape 两倍 hard_sum 才 forced_skipped
            if hard_sum > 0 and total > hard_sum * 2:
                forced_whole_skip = True

            if forced_whole_skip:
                return value, {
                    "action": "forced_skipped",
                    "chars": total,
                    "max_chars": max_chars,
                    "hard_ceiling_total": hard_sum,
                    "reason": "total_exceed_2x_hard_ceiling_sum",
                }
            if per_para_truncated or emphasis_protected:
                return new_paragraphs, {
                    "action": "per_paragraph_truncated",
                    "chars_before": total,
                    "chars_after": sum(_count_chars(p) for p in new_paragraphs
                                        if isinstance(p, str)
                                        and p != "__skip_paragraph__"),
                    "per_paragraph_truncated": per_para_truncated,
                    "emphasis_protected": emphasis_protected,
                    "max_chars": max_chars,
                }
            return new_paragraphs, None

        # 无 per_paragraph：走老逻辑，但用 hard_ceiling 做阈值
        if total <= effective_truncate:
            return value, None
        if total > effective_truncate * 2:
            return value, {
                "action": "forced_skipped",
                "chars": total, "max_chars": max_chars,
                "hard_ceiling": effective_truncate,
                "reason": "total_paragraph_chars_exceed_2x_hard_ceiling",
            }
        ratio = effective_truncate / total
        new_paragraphs = []
        for p in value:
            if not isinstance(p, str):
                new_paragraphs.append(p)
                continue
            target = max(1, int(_count_chars(p) * ratio))
            new_paragraphs.append(_truncate_with_ellipsis(p, target))
        return new_paragraphs, {
            "action": "truncated",
            "chars_before": total,
            "chars_after": sum(_count_chars(p) for p in new_paragraphs
                               if isinstance(p, str)),
            "max_chars": max_chars,
            "hard_ceiling": effective_truncate,
        }

    return value, None


def _truncate_with_ellipsis(text: str, max_chars: int) -> str:
    """截断到约 max_chars 个有效字符（保留空白），末尾加 …。"""
    if max_chars <= 1:
        return ELLIPSIS
    out_chars = []
    visible = 0
    for c in text:
        if visible >= max_chars - 1:
            break
        out_chars.append(c)
        if c not in (" ", "\n", "\t", "\r", "\u3000"):
            visible += 1
    out = "".join(out_chars).rstrip()
    return out + ELLIPSIS


# === 低层 XML 操作 ===

def _clear_to_empty_preserve_format(text_frame) -> None:
    """清空段落文字但绝不删除段落节点和格式节点。"""
    for para in text_frame.paragraphs:
        for r in para.runs:
            r.text = ""
        try:
            p_el = para._p
            for t in p_el.iter(qn("a:t")):
                t.text = ""
        except Exception:
            pass


def _pt_to_sz(pt: float | int | None) -> str | None:
    """PT -> OOXML sz (百分之一 pt 整数字符串)；非法值返回 None。"""
    if pt is None:
        return None
    try:
        v = float(pt)
    except (TypeError, ValueError):
        return None
    if v <= 0:
        return None
    return str(int(round(v * 100)))


def _pin_run_font_size(run_el, pt: float | int | None) -> None:
    """在 <a:r>/<a:rPr> 上显式写入 sz=%.2f*100；已有则覆盖。"""
    sz = _pt_to_sz(pt)
    if sz is None or run_el is None:
        return
    from lxml import etree
    rPr = run_el.find(qn("a:rPr"))
    if rPr is None:
        rPr = etree.SubElement(run_el, qn("a:rPr"))
        # <a:r> schema 要求 rPr 在 t 之前
        run_el.remove(rPr)
        run_el.insert(0, rPr)
    rPr.set("sz", sz)


def _pin_paragraph_def_font_size(p_el, pt: float | int | None) -> None:
    """在 <a:p>/<a:pPr>/<a:defRPr> 上显式写入 sz。

    作用：防止 listStyle 的 lvlXpPr/defRPr 或 txBody 默认字号覆盖新段。
    """
    sz = _pt_to_sz(pt)
    if sz is None or p_el is None:
        return
    from lxml import etree
    pPr = p_el.find(qn("a:pPr"))
    if pPr is None:
        pPr = etree.SubElement(p_el, qn("a:pPr"))
        # pPr 必须是 <a:p> 的第一个子节点
        p_el.remove(pPr)
        p_el.insert(0, pPr)
    defRPr = pPr.find(qn("a:defRPr"))
    if defRPr is None:
        defRPr = etree.SubElement(pPr, qn("a:defRPr"))
    defRPr.set("sz", sz)


def _replace_paragraph_text(paragraph, new_text: str,
                            font_size_pt: float | int | None = None) -> None:
    """保留 <a:pPr> 与首个 <a:r> 的 <a:rPr> 前提下设置段落文字。

    若 font_size_pt 非空：显式 pin 所有 run 的 rPr.sz 与 pPr/defRPr.sz；
    若原段无 run，则创建一个带 rPr.sz 的新 run 写入新文本，杜绝"继承默认字号"。
    """
    from lxml import etree
    runs = list(paragraph.runs)

    # 任意情况都尝试 pin pPr/defRPr/sz，给未来克隆/继承兜底
    _pin_paragraph_def_font_size(paragraph._p, font_size_pt)

    if not runs:
        # 兜底：手动创建带 rPr 的 run
        p_el = paragraph._p
        # 保留首个 <a:pPr>，其余清除（防重复）
        for child in list(p_el):
            if child.tag == qn("a:pPr"):
                continue
            if child.tag == qn("a:endParaRPr"):
                continue
            p_el.remove(child)
        # 插入 run，位置：pPr 之后、endParaRPr 之前
        new_r = etree.SubElement(p_el, qn("a:r"))
        # 先 rPr 后 t
        new_rPr = etree.SubElement(new_r, qn("a:rPr"))
        new_rPr.set("lang", "zh-CN")
        sz = _pt_to_sz(font_size_pt)
        if sz:
            new_rPr.set("sz", sz)
        new_t = etree.SubElement(new_r, qn("a:t"))
        new_t.text = new_text
        # 确保 new_r 在 endParaRPr 之前
        epr = p_el.find(qn("a:endParaRPr"))
        if epr is not None:
            p_el.remove(new_r)
            p_el.insert(list(p_el).index(epr), new_r)
        return

    # 有 run：第一个 run 写入，其他清空；所有 run 的 rPr.sz pin
    runs[0].text = new_text
    _pin_run_font_size(runs[0]._r, font_size_pt)
    for r in runs[1:]:
        r.text = ""
        _pin_run_font_size(r._r, font_size_pt)


def _clone_paragraph_for_extra(text_frame, template_para,
                               font_size_pt: float | int | None = None):
    """克隆模板段落（含 <a:pPr> bullet/缩进/level）作为新段，追加到 text_frame。

    若 font_size_pt 非空，克隆后立即覆盖 pPr/defRPr.sz 与所有 run rPr.sz，
    防止克隆出的段落继承模板段原字号（尤其原段是大字标题时）。
    """
    txBody = text_frame._txBody
    new_p = copy.deepcopy(template_para._p)
    # 清空新段的 <a:t>
    for t in new_p.iter(qn("a:t")):
        t.text = ""
    txBody.append(new_p)
    # 重新枚举 paragraphs，找出最后一个对应 new_p
    paras = list(text_frame.paragraphs)
    new_para = paras[-1]

    if font_size_pt is not None:
        _pin_paragraph_def_font_size(new_para._p, font_size_pt)
        for r in new_para.runs:
            _pin_run_font_size(r._r, font_size_pt)

    return new_para


SKIP_PARAGRAPH_TOKEN = "__skip_paragraph__"


def _font_for_idx(i: int,
                  per_paragraph_font_sizes: list[float] | None,
                  fallback: float | None) -> float | None:
    """取第 i 段目标字号：优先 per_paragraph_font_sizes[i]，否则 fallback。"""
    if per_paragraph_font_sizes and i < len(per_paragraph_font_sizes):
        v = per_paragraph_font_sizes[i]
        if v:
            return float(v)
    return float(fallback) if fallback else None


def _apply_paragraphs(text_frame,
                      values: list[str],
                      per_paragraph_font_sizes: list[float] | None = None,
                      shape_font_size: float | None = None) -> None:
    """多段落写入 + 逐段 pin 字号；段数不足则克隆并继续 pin。

    values[i] == '__skip_paragraph__' 表示保留该段原文不动（emphasis 段保护）。
    """
    paragraphs = list(text_frame.paragraphs)
    if not paragraphs:
        # 完全空 textbox：用 text 一次性写入（跳过保护段）
        joined = "\n".join(v for v in values if v != SKIP_PARAGRAPH_TOKEN)
        text_frame.text = joined
        # 写入后再逐段 pin
        for i, p in enumerate(text_frame.paragraphs):
            fz = _font_for_idx(i, per_paragraph_font_sizes, shape_font_size)
            if fz is not None:
                _pin_paragraph_def_font_size(p._p, fz)
                for r in p.runs:
                    _pin_run_font_size(r._r, fz)
        return

    n_existing = len(paragraphs)
    n_new = len(values)

    # 写入已有段
    for i in range(min(n_existing, n_new)):
        if values[i] == SKIP_PARAGRAPH_TOKEN:
            continue  # 保留原段
        fz = _font_for_idx(i, per_paragraph_font_sizes, shape_font_size)
        _replace_paragraph_text(paragraphs[i], values[i], font_size_pt=fz)

    # 新段超出原段数 → 克隆最后一段的 pPr 追加
    if n_new > n_existing:
        last_template = paragraphs[-1]
        for j in range(n_existing, n_new):
            if values[j] == SKIP_PARAGRAPH_TOKEN:
                continue
            fz = _font_for_idx(j, per_paragraph_font_sizes, shape_font_size)
            new_para = _clone_paragraph_for_extra(text_frame, last_template,
                                                  font_size_pt=fz)
            _replace_paragraph_text(new_para, values[j], font_size_pt=fz)

    # 原段数 > 新段数 → 清空多余段（保留 pPr 节点；不强制 pin 避免破坏原样式）
    elif n_existing > n_new:
        for j in range(n_new, n_existing):
            _replace_paragraph_text(paragraphs[j], "")


def _apply_to_text_frame(text_frame, value,
                         font_size_pt: float | None = None,
                         per_paragraph_font_sizes: list[float] | None = None,
                         ) -> None:
    if isinstance(value, list):
        _apply_paragraphs(text_frame, [str(p) for p in value],
                          per_paragraph_font_sizes=per_paragraph_font_sizes,
                          shape_font_size=font_size_pt)
        return
    paragraphs = list(text_frame.paragraphs)
    if not paragraphs:
        text_frame.text = value if isinstance(value, str) else str(value)
        # 写入后 pin
        for p in text_frame.paragraphs:
            if font_size_pt is not None:
                _pin_paragraph_def_font_size(p._p, font_size_pt)
                for r in p.runs:
                    _pin_run_font_size(r._r, font_size_pt)
        return
    if isinstance(value, str):
        # 单段 shape：优先用 per_paragraph_font_sizes[0]（若有），否则 shape 级
        target_fz = _font_for_idx(0, per_paragraph_font_sizes, font_size_pt)
        _replace_paragraph_text(paragraphs[0], value, font_size_pt=target_fz)
        for p in paragraphs[1:]:
            _replace_paragraph_text(p, "")


# === shape_id → shape 路由 ===

def _walk_all_shapes(slide):
    """yield shape，递归展开 group。"""
    def _walk(container):
        for shape in container.shapes if hasattr(container, "shapes") else container:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _walk(shape)
            else:
                yield shape
    yield from _walk(slide)


def _build_shape_index(prs) -> dict[str, object]:
    """key = slide_{i}::sp_{cNvPr.id} → shape；
       table cell key = slide_{i}::sp_{tableId}::cell_{r}_{c} → cell.text_frame
    """
    idx: dict[str, object] = {}
    for i, slide in enumerate(prs.slides, start=1):
        for shape in _walk_all_shapes(slide):
            try:
                cnv = shape.shape_id
            except Exception:
                continue
            if cnv is None:
                continue
            base_key = f"slide_{i}::sp_{cnv}"
            idx[base_key] = shape
            # 表格：注册每个 cell
            try:
                if shape.has_table:
                    for r_idx, row in enumerate(shape.table.rows):
                        for c_idx, cell in enumerate(row.cells):
                            cell_key = f"{base_key}::cell_{r_idx}_{c_idx}"
                            idx[cell_key] = cell.text_frame
            except Exception:
                pass
    return idx


# === mapping 结构适配 ===

def _flatten_mapping(mapping: dict) -> tuple[dict, dict]:
    flat: dict[str, object] = {}
    meta: dict[str, dict] = {}

    if isinstance(mapping, dict) and "slides" in mapping \
            and isinstance(mapping["slides"], dict):
        for _idx_str, page in mapping["slides"].items():
            for a in page.get("assignments", []):
                key = a["key"]
                flat[key] = a["value"]
                meta[key] = {
                    "shape_id": a.get("shape_id"),
                    "role": a.get("role"),
                    "source": a.get("source"),
                    "preserve_action": a.get("preserve_action"),
                    "char_limit": a.get("char_limit"),
                    "hard_ceiling_chars": a.get("hard_ceiling_chars"),
                    "per_paragraph_char_limits": a.get(
                        "per_paragraph_char_limits"),
                    "is_placeholder_text": a.get("is_placeholder_text"),
                    "font_size_pt": a.get("font_size_pt"),
                    "per_paragraph_font_size_pt": a.get(
                        "per_paragraph_font_size_pt") or [],
                }
        return flat, meta

    # 兼容老扁平 schema
    for k, v in mapping.items():
        if k in ("schema_version", "pending_shrink_requests", "meta",
                 "fallback_required", "global_style_guide"):
            continue
        flat[k] = v
        meta[k] = {}
    return flat, meta


def _build_story_index(story: dict | None) -> dict[str, dict]:
    """key = shape_id → {origin, role, preserve_action, current_text, char_limit}。"""
    if not story:
        return {}
    idx: dict[str, dict] = {}
    for s in story.get("slides", []):
        for cs in s.get("content_shapes", []):
            idx[cs["shape_id"]] = {
                "origin": "content",
                "role": "content_shape",
                "preserve_action": None,
                "current_text": cs.get("original_text", ""),
                "char_limit": cs.get("char_limit"),
                "hard_ceiling_chars": cs.get("hard_ceiling_chars"),
                "per_paragraph": cs.get("per_paragraph"),
                "is_placeholder_text": bool(cs.get("is_placeholder_text")),
                "font_size_pt": cs.get("font_size_pt"),
                "per_paragraph_font_size_pt":
                    cs.get("per_paragraph_font_size_pt") or [],
            }
        for nc in s.get("non_content_shapes", []):
            idx[nc["shape_id"]] = {
                "origin": "non_content",
                "role": nc.get("role"),
                "preserve_action": nc.get("preserve_action", "keep_original"),
                "current_text": nc.get("current_text", ""),
                "char_limit": None,
                "is_placeholder_text": False,
            }
    return idx


# === apply ===

def apply(pptx_in: Path, mapping: dict, pptx_out: Path, *,
          story: dict | None = None,
          force_apply: bool = False) -> dict:
    pptx_out.parent.mkdir(parents=True, exist_ok=True)
    if pptx_in.resolve() != pptx_out.resolve():
        shutil.copyfile(pptx_in, pptx_out)

    # capacity gate
    pending = mapping.get("pending_shrink_requests") or []
    if pending and not force_apply:
        return {
            "ok": False,
            "error": "capacity_mismatch_gate",
            "message": "pending_shrink_requests 非空；请先跑 Prompt D shrink 后再 apply",
            "pending_count": len(pending),
            "samples": pending[:5],
        }

    prs = Presentation(str(pptx_out))
    flat_mapping, meta_index = _flatten_mapping(mapping)
    story_idx = _build_story_index(story)
    shape_idx = _build_shape_index(prs)

    applied: list[str] = []
    cleared: list[str] = []
    skipped: list[str] = []
    forced_skipped: list[dict] = []
    truncated: list[dict] = []
    per_paragraph_truncated_count = 0
    emphasis_protected_count = 0
    font_size_pinned_count = 0
    used_keys: set[str] = set()

    for key, value in flat_mapping.items():
        target = shape_idx.get(key)
        if target is None:
            continue
        used_keys.add(key)

        story_info = story_idx.get(key, {})
        a_meta = meta_index.get(key, {})
        char_limit = a_meta.get("char_limit") or story_info.get("char_limit")
        hard_ceiling = a_meta.get("hard_ceiling_chars") or \
            story_info.get("hard_ceiling_chars")
        per_paragraph_limits = a_meta.get("per_paragraph_char_limits") or \
            story_info.get("per_paragraph")
        # 字号 pin：shape 级 + 段级
        shape_font_pt = a_meta.get("font_size_pt") or \
            story_info.get("font_size_pt")
        per_para_font_pt = a_meta.get("per_paragraph_font_size_pt") or \
            story_info.get("per_paragraph_font_size_pt") or []
        # 若 a_meta 的 per_paragraph_char_limits 里带 font_size_pt，也采一份（兼容）
        if not per_para_font_pt and per_paragraph_limits:
            extracted = [pp.get("font_size_pt") for pp in per_paragraph_limits
                         if isinstance(pp, dict)]
            if any(extracted):
                per_para_font_pt = extracted

        # non_content + keep_original：只允许 __skip__/__clear__；其他强制跳过
        if story_info.get("origin") == "non_content":
            if story_info.get("preserve_action") == "keep_original" \
                    and value not in ("__skip__", "__clear__"):
                forced_skipped.append({
                    "key": key, "reason": "non_content_keep_original",
                    "role": story_info.get("role"),
                    "offending_value": (value if isinstance(value, str)
                                        else f"list[{len(value)}]")[:40],
                })
                skipped.append(key)
                continue

        if value == "__skip__":
            skipped.append(key)
            continue

        # 取 text_frame：shape 或直接 cell.text_frame
        if hasattr(target, "text_frame"):
            try:
                tf = target.text_frame
            except Exception:
                skipped.append(key)
                continue
        elif hasattr(target, "paragraphs"):
            tf = target  # 已经是 text_frame
        else:
            skipped.append(key)
            continue

        if value == "__clear__":
            try:
                _clear_to_empty_preserve_format(tf)
                cleared.append(key)
            except Exception as e:  # noqa: BLE001
                skipped.append(f"{key} (clear error: {e})")
            continue

        # L5.1 装饰保护兜底
        original_text = story_info.get("current_text") or ""
        if not original_text:
            try:
                original_text = tf.text or ""
            except Exception:
                original_text = ""
        if _decoration_protection_triggered(original_text, value):
            forced_skipped.append({
                "key": key,
                "reason": "decoration_text_protected",
                "decoration_subtype": _looks_like_decoration(original_text),
                "original_text": original_text.strip()[:20],
                "offending_value": (value if isinstance(value, str)
                                    else f"list[{len(value)}]")[:40],
            })
            skipped.append(key)
            continue

        # L5.2 硬字数闸门
        new_value, gate_meta = _enforce_char_limit(
            value, char_limit,
            per_paragraph=per_paragraph_limits,
            hard_ceiling=hard_ceiling,
        )
        if gate_meta and gate_meta.get("action") == "forced_skipped":
            forced_skipped.append({
                "key": key,
                "reason": gate_meta.get("reason", "char_limit_exceeded"),
                "char_limit": char_limit,
                "chars": gate_meta.get("chars"),
                "offending_value": (value if isinstance(value, str)
                                    else f"list[{len(value)}]")[:40],
            })
            skipped.append(key)
            continue
        if gate_meta and gate_meta.get("action") == "truncated":
            truncated.append({
                "key": key,
                "chars_before": gate_meta.get("chars_before"),
                "chars_after": gate_meta.get("chars_after"),
                "max_chars": gate_meta.get("max_chars"),
            })
        if gate_meta and gate_meta.get("action") == "per_paragraph_truncated":
            per_paragraph_truncated_count += gate_meta.get(
                "per_paragraph_truncated", 0)
            emphasis_protected_count += gate_meta.get(
                "emphasis_protected", 0)
            truncated.append({
                "key": key,
                "chars_before": gate_meta.get("chars_before"),
                "chars_after": gate_meta.get("chars_after"),
                "per_paragraph_truncated": gate_meta.get(
                    "per_paragraph_truncated"),
                "emphasis_protected": gate_meta.get("emphasis_protected"),
            })

        try:
            _apply_to_text_frame(tf, new_value,
                                 font_size_pt=shape_font_pt,
                                 per_paragraph_font_sizes=per_para_font_pt)
            applied.append(key)
            if shape_font_pt or any(per_para_font_pt):
                font_size_pinned_count += 1
        except Exception as e:  # noqa: BLE001
            skipped.append(f"{key} (error: {e})")

    not_found = [k for k in flat_mapping if k not in used_keys]
    prs.save(str(pptx_out))

    return {
        "ok": True,
        "out": str(pptx_out),
        "applied": len(applied),
        "cleared": len(cleared),
        "skipped": len(skipped),
        "truncated_count": len(truncated),
        "truncated_samples": truncated[:10],
        "per_paragraph_truncated_count": per_paragraph_truncated_count,
        "emphasis_protected_count": emphasis_protected_count,
        "font_size_pinned_count": font_size_pinned_count,
        "forced_skipped_count": len(forced_skipped),
        "forced_skipped": forced_skipped[:20],
        "not_found": not_found[:20],
        "not_found_total": len(not_found),
    }


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("pptx", type=Path)
    ap.add_argument("--mapping", required=True, type=Path)
    ap.add_argument("--out", required=True, type=Path)
    ap.add_argument("--story", type=Path, default=None,
                    help="可选 template_story.json；开启 preserve_action 兜底")
    ap.add_argument("--force-apply", action="store_true",
                    help="即使 pending_shrink_requests 非空也落盘"
                         "（仍受 char_limit 闸门和装饰保护约束）")
    args = ap.parse_args()

    if not args.pptx.exists():
        print(f"ERROR: pptx not found: {args.pptx}", file=sys.stderr)
        return 1
    if not args.mapping.exists():
        print(f"ERROR: mapping not found: {args.mapping}", file=sys.stderr)
        return 1

    mapping = json.loads(args.mapping.read_text(encoding="utf-8"))
    story = None
    if args.story and args.story.exists():
        story = json.loads(args.story.read_text(encoding="utf-8"))

    result = apply(args.pptx, mapping, args.out, story=story,
                   force_apply=args.force_apply)
    print(json.dumps(result, ensure_ascii=False))
    return 0 if result.get("ok") else 2


if __name__ == "__main__":
    sys.exit(main())
