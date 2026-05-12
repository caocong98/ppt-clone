"""把 PPT 模板解析为 **shape 级字符驱动** 的叙事骨架（template_story.json，schema v2）。

三层混合架构：
- L1 语义层：story_role + 装饰过滤（L1/L2/L3）+ is_placeholder_text + global_style_guide
- L2 结构提示层：shape_group 弱提示（字数桶 + bbox 网格对齐）—— 仅 LLM 可见，不影响 apply
- L3 硬约束层：每个 content_shape 独立的 char_limit + paragraph_count + font_size + estimated_capacity

非内容判定（保留旧规则）：
- L1 规则层：placeholder 类型 / 正则 / 位置 / 面积字号 / 占位词典
- L2 跨页重复层：跨幻灯片相同文本 + 位置聚类 → logo_text / running_header/footer
- L3 视觉层：L1+L2 未命中的短文/角落/badge 进 vision_ambiguous[]，由 Prompt B 判定

CLI:
  python parse_template_story.py <pptx> --out template_story.json
      [--logo-action keep_original|clear_to_empty]
      [--shape-group-min 3] [--no-debug-dump]

  python parse_template_story.py <pptx> --out template_story.json \
      --resume vision_result.json    # 回灌 Prompt B 的判定结果
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

SCHEMA_VERSION = 2

DEFAULT_THRESHOLDS = {
    # 跨页/装饰
    "cross_slide_repeat_min": 3,
    "logo_corner_distance_pct": 5.0,
    "decoration_micro_area_pct": 0.4,
    "decoration_micro_font_pt": 12.0,
    "decoration_micro_text_len": 4,
    # 视觉可疑
    "vision_ambiguous_text_len_max": 12,
    "vision_ambiguous_area_pct_max": 2.0,
    "vision_ambiguous_edge_distance_pct": 8.0,
    # shape_group 弱提示
    "shape_group_min_members": 3,
    "shape_group_alignment_var_max_pct": 1.5,
    "char_bucket_edges": [10, 30, 80],
    # 字符容量估算（收紧 + 留白）
    "char_capacity_safety_factor": 0.85,
    "char_capacity_line_height_factor": 1.4,
    "text_inset_pt": 3.6,                      # bodyPr 缺省值，实测读取见 _read_body_inset
    "chinese_char_width_factor": 1.1,          # 中文字形实际宽度相对字号（含字间距）
    "char_limit_fill_ratio": 0.65,             # 占位 shape 推荐填充比例（留白 35%）
    "emphasis_font_ratio": 1.25,               # 段字号 >= 中位数 * 该比例 → is_emphasis
    "char_limit_max_multiplier_lower": 1.3,
    "char_limit_max_multiplier_upper": 1.1,
    "char_limit_min_multiplier": 0.6,
}

# === 文本规则库 ===
_NUMBER_MARKER_RE = re.compile(r"^\s*0?\d{1,2}\s*[\.．、:：\-/]?\s*$")
_ENUM_PREFIX_RE = re.compile(r"^\s*(0?[1-9]|1[0-9])\s*[\.．、\s:：]\s*")

# === 序号装饰字符集（覆盖汉字大写/罗马/圆圈/ASCII 序号）===
_CJK_NUMERAL_CHARS = set("壹贰叁肆伍陆柒捌玖拾佰仟一二三四五六七八九十零百千万")
_ROMAN_NUMERAL_RE = re.compile(
    r"^\s*[ⅠⅡⅢⅣⅤⅥⅦⅧⅨⅩⅪⅫⅰⅱⅲⅳⅴⅵⅶⅷⅸⅹ]{1,4}\s*[\.．、]?\s*$")
_CIRCLED_DIGIT_RE = re.compile(
    r"^\s*[①②③④⑤⑥⑦⑧⑨⑩⑪⑫⑬⑭⑮⑯⑰⑱⑲⑳❶❷❸❹❺❻❼❽❾❿]\s*$")


def _is_ordinal_decoration(text: str) -> str | None:
    """装饰序号识别：返回命中的子类型，未命中返回 None。"""
    if text is None:
        return None
    t = text.strip()
    if not t:
        return None
    if _NUMBER_MARKER_RE.match(t):
        return "ordinal_arabic"
    if _ROMAN_NUMERAL_RE.match(t):
        return "ordinal_roman"
    if _CIRCLED_DIGIT_RE.match(t):
        return "ordinal_circled"
    if len(t) <= 2 and all(c in _CJK_NUMERAL_CHARS for c in t):
        return "ordinal_cjk"
    return None


_LOGO_TEXT_KEYWORDS = (
    "logo", "your logo", "your slogan", "company name", "brand", "trademark",
    "®", "©", "™", "wps", "office", "powerpoint", "ppt",
    "docer", "稻壳", "金山", "wpsoffice", "微软",
    "click to add", "click here", "添加标题", "请添加", "在此添加",
    "your text here",
)

_LOGO_DEFAULT_ACTION = "keep_original"


_CONTENT_PLACEHOLDER_PATTERNS = [
    re.compile(r"please enter", re.IGNORECASE),
    re.compile(r"enter your (content|text)", re.IGNORECASE),
    re.compile(r"paste here", re.IGNORECASE),
    re.compile(r"type (here|your)", re.IGNORECASE),
    re.compile(r"insert (text|content)", re.IGNORECASE),
    re.compile(r"add your (text|content)", re.IGNORECASE),
    re.compile(r"click here to add", re.IGNORECASE),
    re.compile(r"(\u8f93\u5165|\u586b\u5199)(\u5185\u5bb9|\u6587[\u672c\u5b57])"),
    re.compile(r"\u6b64\u5904\u8f93\u5165"),
    re.compile(r"\u8bf7\u8f93\u5165"),
    re.compile(r"\u5728\u6b64(\u8f93\u5165|\u6dfb\u52a0|\u7f16\u8f91)"),
]

_GENERIC_PLACEHOLDER_KEYWORDS = {
    "内容概述", "内容描述", "内容标题", "内容说明",
    "存在问题", "解决方案", "工作计划", "工作内容", "工作总结",
    "项目概述", "项目描述", "项目总结", "项目进展",
    "数据分析", "数据概述", "重点项目",
    "标题", "副标题", "小标题", "正文",
}


def _is_content_placeholder(text: str) -> bool:
    """识别"内容占位提示语"——这种文本应被新内容替换，不应被装饰过滤吃掉。"""
    if text is None:
        return False
    t = text.strip()
    if not t:
        return False
    for pat in _CONTENT_PLACEHOLDER_PATTERNS:
        if pat.search(t):
            return True
    return t in _GENERIC_PLACEHOLDER_KEYWORDS


_STYLE_TAG_HINTS_RE = re.compile(
    r"^[A-Z][A-Z0-9 \-&'.]{2,29}$"  # ALL CAPS short
)


def _is_latin_decoration_tag(text: str, bbox: dict, ph_type: str | None) -> bool:
    """识别装饰性的英文样式标签（如 QINGMING FESTIVAL / DESIGN）。"""
    if ph_type in ("TITLE", "CENTER_TITLE", "SUBTITLE"):
        return False
    if not text:
        return False
    t = text.strip()
    if not t:
        return False
    if len(t) > 30:
        return False
    if not all(ord(c) < 128 for c in t):
        return False
    if any(c.isdigit() and int(c) >= 0 for c in t if c.isdigit()):
        if len(t) <= 4:
            return False
    if not _STYLE_TAG_HINTS_RE.match(t):
        return False
    edge = _distance_to_edge_pct(bbox)
    small = bbox.get("area_pct", 0) < 1.5
    return edge < 8.0 or small


# === 工具 ===

EMU_PER_PT = 12700


def _emu_to_pt(emu: int) -> float:
    return emu / EMU_PER_PT


def _is_chinese(s: str) -> bool:
    return any("\u4e00" <= c <= "\u9fff" for c in (s or ""))


def _normalize_text(s: str) -> str:
    return re.sub(r"\s+", " ", (s or "").strip()).lower()


def _bbox_pct(shape, slide_w_emu: int, slide_h_emu: int) -> dict:
    try:
        left = shape.left or 0
        top = shape.top or 0
        w = shape.width or 0
        h = shape.height or 0
    except Exception:
        return {"left": 0.0, "top": 0.0, "w": 0.0, "h": 0.0, "area_pct": 0.0,
                "left_emu": 0, "top_emu": 0, "w_emu": 0, "h_emu": 0}
    sw = max(1, slide_w_emu)
    sh = max(1, slide_h_emu)
    return {
        "left": round(left / sw * 100, 2),
        "top": round(top / sh * 100, 2),
        "w": round(w / sw * 100, 2),
        "h": round(h / sh * 100, 2),
        "area_pct": round(w * h / (sw * sh) * 100, 4),
        "left_emu": left,
        "top_emu": top,
        "w_emu": w,
        "h_emu": h,
    }


def _distance_to_edge_pct(bbox: dict) -> float:
    cx = bbox["left"] + bbox["w"] / 2
    cy = bbox["top"] + bbox["h"] / 2
    return min(cx, 100 - cx, cy, 100 - cy)


def _bbox_region(bbox: dict) -> str:
    """返回如 'top-left' / 'mid-center' / 'bottom-right' 的位置标签。"""
    cx = bbox.get("left", 0) + bbox.get("w", 0) / 2
    cy = bbox.get("top", 0) + bbox.get("h", 0) / 2
    if cy < 33:
        v = "top"
    elif cy < 67:
        v = "mid"
    else:
        v = "bottom"
    if cx < 33:
        h = "left"
    elif cx < 67:
        h = "center"
    else:
        h = "right"
    return f"{v}-{h}"


def _walk_shapes(container, path: str = ""):
    """yield (shape, group_path)，递归展开 group。"""
    for shape in container.shapes if hasattr(container, "shapes") else container:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            sub = f"{path}/{shape.name}" if path else shape.name
            yield from _walk_shapes(shape, sub)
        else:
            yield shape, path


def _resolve_shape_id(shape, slide_idx: int) -> str:
    """全局唯一 shape_id：slide_{i}::sp_{cNvPr.id}。

    若 cNvPr.id 不可得（极少数特殊 shape），降级为 hash(name+bbox)。
    """
    try:
        cnv_id = shape.shape_id
        if cnv_id is not None:
            return f"slide_{slide_idx}::sp_{cnv_id}"
    except Exception:
        pass
    name = getattr(shape, "name", "unknown") or "unknown"
    return f"slide_{slide_idx}::name_{name}"


def _resolve_font_size(run=None, paragraph=None, text_frame=None,
                       shape=None, *, default: float = 18.0) -> float:
    """字号回退链：run → paragraph defRPr → text_frame lstStyle → shape defaults → default。

    返回 pt。无法确定时返回 default（默认 18.0）。
    """
    # 1) run.font.size
    try:
        if run is not None and run.font.size is not None:
            return float(run.font.size.pt)
    except Exception:
        pass
    # 2) paragraph 级 defRPr
    try:
        if paragraph is not None:
            pPr = getattr(paragraph, "_pPr", None)
            if pPr is not None:
                defRPr = pPr.find(qn("a:defRPr"))
                if defRPr is not None and defRPr.get("sz") is not None:
                    sz_val = int(defRPr.get("sz"))
                    return sz_val / 100.0
    except Exception:
        pass
    # 3) text_frame 的 lstStyle 各层级
    try:
        if text_frame is not None:
            tf_el = text_frame._txBody
            lstStyle = tf_el.find(qn("a:lstStyle"))
            if lstStyle is not None:
                for tag in ("lvl1pPr", "lvl2pPr", "lvl3pPr", "defPPr"):
                    pNode = lstStyle.find(qn(f"a:{tag}"))
                    if pNode is not None:
                        defRPr = pNode.find(qn("a:defRPr"))
                        if defRPr is not None and defRPr.get("sz") is not None:
                            sz_val = int(defRPr.get("sz"))
                            return sz_val / 100.0
    except Exception:
        pass
    # 4) shape 的 sp_pr 默认（常无效，留作钩子）
    return default


def _read_body_inset(text_frame) -> dict:
    """从 <a:bodyPr> 读实际 lIns/tIns/rIns/bIns，EMU → pt (1 pt = 12700 EMU)。

    PPTX 默认：lIns/rIns=91440 EMU (7.2pt)，tIns/bIns=45720 EMU (3.6pt)。
    bodyPr 属性缺失时不回传（由 _resolve_insets 用 threshold 默认值填补）。
    """
    out: dict = {}
    if text_frame is None:
        return out
    try:
        body_pr = text_frame._txBody.find(qn("a:bodyPr"))
    except Exception:
        return out
    if body_pr is None:
        return out
    for attr, key in (("lIns", "lIns_pt"), ("tIns", "tIns_pt"),
                      ("rIns", "rIns_pt"), ("bIns", "bIns_pt")):
        v = body_pr.get(attr)
        if v is not None:
            try:
                out[key] = int(v) / 12700.0
            except (TypeError, ValueError):
                continue
    # 未显式给时回退 PPTX 官方默认值
    out.setdefault("lIns_pt", 7.2)
    out.setdefault("rIns_pt", 7.2)
    out.setdefault("tIns_pt", 3.6)
    out.setdefault("bIns_pt", 3.6)
    return out


def _gather_text(text_frame) -> tuple[str, list[str], float, list[float]]:
    """返回 (full_text, paragraphs, dominant_font_size_pt, per_paragraph_dominant_size)。

    每段字号采用 "dominant run" 策略：取段内最长 run 的字号；若段内所有 run 都
    为空则退回 _resolve_font_size(paragraph=...)。这样大字号标题 + 小字号注释
    混排时，新文本不会继承错误的字号。
    """
    paragraphs: list[str] = []
    per_para_sizes: list[float] = []

    for p in text_frame.paragraphs:
        text = "".join(r.text or "" for r in p.runs)
        if not text and p.text:
            text = p.text
        paragraphs.append(text)

        # 每段取"最长 run"的字号作为 dominant
        run_candidates: list[tuple[int, float]] = []
        for r in p.runs:
            t = r.text or ""
            if not t.strip():
                continue
            sz = _resolve_font_size(run=r, paragraph=p, text_frame=text_frame)
            run_candidates.append((len(t), sz))

        if run_candidates:
            dominant = max(run_candidates, key=lambda x: x[0])[1]
        else:
            dominant = _resolve_font_size(paragraph=p, text_frame=text_frame)
        per_para_sizes.append(float(dominant))

    full_text = "\n".join(paragraphs)
    # 整 shape 的"主导字号"：取段字号按字数加权；无字则取第一段
    if paragraphs and any(paragraphs):
        total_w = 0.0
        total_n = 0
        for i, para_text in enumerate(paragraphs):
            n = len(para_text)
            if n > 0:
                total_w += per_para_sizes[i] * n
                total_n += n
        weighted = (total_w / total_n) if total_n > 0 else per_para_sizes[0]
    else:
        weighted = (per_para_sizes[0] if per_para_sizes
                    else (_resolve_font_size(text_frame=text_frame)
                          if text_frame else 18.0))
    return full_text, paragraphs, weighted, per_para_sizes


def _paragraph_has_bullet(paragraph) -> bool:
    try:
        pPr = paragraph._pPr
        if pPr is None:
            return False
        for tag in ("buChar", "buAutoNum", "buBlip"):
            if pPr.find(qn(f"a:{tag}")) is not None:
                return True
    except Exception:
        pass
    t = paragraph.text or ""
    t_stripped = t.strip()
    if t_stripped and t_stripped[0] in "•●◆▪·●○-–—*":
        return True
    return False


def _text_frame_has_any_bullet(text_frame) -> bool:
    try:
        for p in text_frame.paragraphs:
            if _paragraph_has_bullet(p):
                return True
    except Exception:
        pass
    return False


def _count_chars(text: str) -> int:
    """可见 Unicode code point 计数（中英都按 1，剔除空白控制字符）。"""
    if not text:
        return 0
    return sum(1 for c in text if c not in (" ", "\n", "\t", "\r", "\u3000"))


def _resolve_insets(thresholds: dict, body_insets: dict | None) -> tuple[float, float, float, float]:
    """返回 (lIns, tIns, rIns, bIns) pt。bodyPr 未给则用 threshold 默认值。"""
    default_inset = float(thresholds.get("text_inset_pt", 3.6))
    if not body_insets:
        return default_inset, default_inset, default_inset, default_inset
    lins = body_insets.get("lIns_pt")
    tins = body_insets.get("tIns_pt")
    rins = body_insets.get("rIns_pt")
    bins = body_insets.get("bIns_pt")
    return (
        float(lins) if lins is not None else default_inset,
        float(tins) if tins is not None else default_inset,
        float(rins) if rins is not None else default_inset,
        float(bins) if bins is not None else default_inset,
    )


def _estimate_char_capacity(bbox: dict, font_size_pt: float,
                            paragraph_count: int,
                            thresholds: dict,
                            body_insets: dict | None = None) -> int:
    """估算文本框可容纳的字符数（收紧版，扣 inset + 中文系数 + 1.4 行高 + 0.85 safety）。

    公式：
      usable_w = w_pt - lIns - rIns
      usable_h = h_pt - tIns - bIns
      chars_per_line = usable_w / (font_size_pt * chinese_char_width_factor)
      lines          = usable_h / (font_size_pt * line_height_factor)
      capacity       = chars_per_line * lines * safety_factor

    body_insets 若为 None 则回退 threshold 的默认 text_inset_pt。
    """
    if font_size_pt <= 0:
        font_size_pt = 18.0
    w_pt = _emu_to_pt(bbox.get("w_emu", 0))
    h_pt = _emu_to_pt(bbox.get("h_emu", 0))
    if w_pt <= 0 or h_pt <= 0:
        return max(1, paragraph_count)

    lins, tins, rins, bins = _resolve_insets(thresholds, body_insets)
    cn_factor = float(thresholds.get("chinese_char_width_factor", 1.1))
    line_h = float(thresholds.get("char_capacity_line_height_factor", 1.4))
    safety = float(thresholds.get("char_capacity_safety_factor", 0.85))

    usable_w = max(1.0, w_pt - lins - rins)
    usable_h = max(1.0, h_pt - tins - bins)
    chars_per_line = max(1.0, usable_w / (font_size_pt * cn_factor))
    lines = max(1.0, usable_h / (font_size_pt * line_h))
    if paragraph_count > 0:
        lines = max(lines, float(paragraph_count))
    capacity = int(chars_per_line * lines * safety)
    return max(1, capacity)


def _compute_per_paragraph_limits(per_para_sizes: list[float],
                                  per_para_char_counts: list[int],
                                  thresholds: dict) -> list[dict]:
    """为多段 shape 逐段计算 char_limit / is_emphasis / font_size_pt（零估算版）。

    核心原则：char_limit 严格等于原文每段字数；bbox 不再参与。
      - char_limit.max = original_chars_i（原文字数就是设计师留给这段的容量）
      - char_limit.min = max(1, original_chars_i - 2)
      - hard_ceiling_chars = int(original_chars_i * 1.15)
      - is_emphasis = size > median 且 size >= median * emphasis_font_ratio
        （仅作为生成器"短词优先"hint，不再影响 char_limit）
    """
    if not per_para_sizes:
        return []
    # char_counts 长度必须与 sizes 对齐；不够补 0
    counts = list(per_para_char_counts or [])
    while len(counts) < len(per_para_sizes):
        counts.append(0)

    emph_ratio = float(thresholds.get("emphasis_font_ratio", 1.25))

    sorted_sizes = sorted(per_para_sizes)
    n = len(sorted_sizes)
    median = sorted_sizes[n // 2] if n % 2 else (
        (sorted_sizes[n // 2 - 1] + sorted_sizes[n // 2]) / 2.0)

    out: list[dict] = []
    for idx, sz in enumerate(per_para_sizes):
        if sz <= 0:
            sz = 18.0
        oc = max(0, int(counts[idx]))
        # 零字段兜底 1，避免 apply 层以 0 截断
        soft_max = max(1, oc)
        soft_min = max(1, oc - 2) if oc >= 3 else max(1, oc)
        hard_ceiling = max(1, int(soft_max * 1.15))
        is_emphasis = sz > median and sz >= median * emph_ratio

        out.append({
            "idx": idx,
            "font_size_pt": round(float(sz), 2),
            "original_char_count": oc,
            "char_limit": {"min": soft_min, "max": soft_max},
            "hard_ceiling_chars": hard_ceiling,
            "is_emphasis": bool(is_emphasis),
        })
    return out


# === L1 规则层：非内容判定 ===

def _classify_non_content_l1(shape, ph_type: str | None, text: str, bbox: dict,
                             font_size_pt: float, thresholds: dict,
                             *, logo_action: str | None = None,
                             ) -> dict | None:
    """L1 规则层命中则返回 {role, preserve_action, detection}；未命中返回 None。"""
    t_stripped = (text or "").strip()
    effective_logo_action = logo_action or _LOGO_DEFAULT_ACTION

    # 序号装饰
    ordinal_kind = _is_ordinal_decoration(t_stripped)
    if ordinal_kind:
        return {
            "role": "decoration_number",
            "preserve_action": "keep_original",
            "detection": {
                "layer": "L1",
                "signal": f"ordinal:{ordinal_kind}",
                "confidence": 0.95,
            },
        }

    # Latin 装饰 tag
    if _is_latin_decoration_tag(t_stripped, bbox, ph_type):
        return {
            "role": "style_tag",
            "preserve_action": "keep_original",
            "detection": {
                "layer": "L1",
                "signal": "latin_decoration_tag",
                "confidence": 0.85,
            },
        }

    # LOGO 关键词
    low = t_stripped.lower()
    if any(kw in low for kw in _LOGO_TEXT_KEYWORDS):
        # 但如果是通用占位关键词 → 标记为 placeholder 后续处理
        if t_stripped in _GENERIC_PLACEHOLDER_KEYWORDS or _is_content_placeholder(t_stripped):
            return None
        return {
            "role": "logo_text",
            "preserve_action": effective_logo_action,
            "detection": {
                "layer": "L1",
                "signal": "logo_keyword",
                "confidence": 0.9,
            },
        }

    # 微小装饰：小面积 + 小字号 + 短文本
    if (bbox.get("area_pct", 0) < thresholds["decoration_micro_area_pct"]
            and font_size_pt > 0
            and font_size_pt < thresholds["decoration_micro_font_pt"]
            and len(t_stripped) <= thresholds["decoration_micro_text_len"]):
        return {
            "role": "decoration_micro",
            "preserve_action": "keep_original",
            "detection": {
                "layer": "L1",
                "signal": "small_area_small_font_short_text",
                "confidence": 0.8,
                "area_pct": bbox.get("area_pct", 0),
                "font_size_pt": font_size_pt,
            },
        }

    return None


def _classify_special_placeholder(ph_type: str | None) -> dict | None:
    """识别 footer/slide_number/datetime 等模板继承占位符 → 默认 keep_original。"""
    if ph_type == "FOOTER":
        return {
            "role": "footer_placeholder",
            "preserve_action": "keep_original",
            "detection": {"layer": "L1", "signal": "ph_type=FOOTER",
                          "confidence": 1.0},
        }
    if ph_type == "SLIDE_NUMBER":
        return {
            "role": "slide_number_placeholder",
            "preserve_action": "keep_original",
            "detection": {"layer": "L1", "signal": "ph_type=SLIDE_NUMBER",
                          "confidence": 1.0},
        }
    if ph_type == "DATE":
        return {
            "role": "date_placeholder",
            "preserve_action": "keep_original",
            "detection": {"layer": "L1", "signal": "ph_type=DATE",
                          "confidence": 1.0},
        }
    return None


def _classify_graphic_shape(shape, ph_type: str | None, bbox: dict,
                            thresholds: dict) -> dict | None:
    """图表/SmartArt/Picture 整体归 non_content（table 不在此处理，单独拆 cell）。"""
    # Chart
    try:
        if shape.has_chart:
            return {"role": "chart_container", "preserve_action": "keep_original",
                    "detection": {"layer": "L1", "signal": "has_chart",
                                  "confidence": 1.0}}
    except Exception:
        pass
    # SmartArt
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and ph_type == "ORG_CHART":
            return {"role": "smartart_container", "preserve_action": "keep_original",
                    "detection": {"layer": "L1", "signal": "ph_type=ORG_CHART",
                                  "confidence": 1.0}}
        el = getattr(shape, "_element", None)
        if el is not None:
            for g in el.iter(qn("a:graphicData")):
                uri = g.get("uri", "")
                if "diagram" in uri.lower() or "smartArt" in uri:
                    return {"role": "smartart_container",
                            "preserve_action": "keep_original",
                            "detection": {"layer": "L1", "signal": "graphicData:diagram",
                                          "confidence": 1.0}}
    except Exception:
        pass
    # Picture
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or ph_type in ("PICTURE", "BITMAP"):
        edge_d = _distance_to_edge_pct(bbox)
        is_corner = edge_d < thresholds["logo_corner_distance_pct"]
        if is_corner and bbox["area_pct"] < 5.0:
            return {"role": "logo_image", "preserve_action": "keep_original",
                    "detection": {"layer": "L1",
                                  "signal": "picture_in_corner_small_area",
                                  "confidence": 0.85,
                                  "edge_distance_pct": edge_d}}
        return {"role": "decorative_picture", "preserve_action": "keep_original",
                "detection": {"layer": "L1", "signal": "picture_non_logo",
                              "confidence": 0.7}}
    # Media
    if shape.shape_type == MSO_SHAPE_TYPE.MEDIA or ph_type == "MEDIA":
        return {"role": "decorative_picture", "preserve_action": "keep_original",
                "detection": {"layer": "L1", "signal": "media",
                              "confidence": 0.9}}
    return None


# === L3 视觉可疑判定 ===

def _is_vision_ambiguous(text: str, bbox: dict, ph_type: str | None,
                         thresholds: dict) -> list[str] | None:
    t = (text or "").strip()
    if not t:
        return None
    if ph_type in ("TITLE", "CENTER_TITLE", "SUBTITLE"):
        return None
    if bbox.get("area_pct", 0) >= 5.0:
        return None

    signals: list[str] = []
    tlen = len(t)
    if tlen <= thresholds["vision_ambiguous_text_len_max"]:
        signals.append(f"short_text(<={thresholds['vision_ambiguous_text_len_max']})")
    area = bbox.get("area_pct", 0)
    if area < thresholds["vision_ambiguous_area_pct_max"]:
        signals.append(f"small_area(<{thresholds['vision_ambiguous_area_pct_max']}%)")
    edge = _distance_to_edge_pct(bbox)
    if edge < thresholds["vision_ambiguous_edge_distance_pct"]:
        signals.append(f"near_edge(<{thresholds['vision_ambiguous_edge_distance_pct']}%)")
    style_hit = bool(_STYLE_TAG_HINTS_RE.match(t))
    if style_hit:
        signals.append("style_tag_pattern")

    if style_hit or len(signals) >= 2:
        return signals
    return None


# === L2 跨页重复分析 ===

def _cross_slide_repeat_analyze(all_texts: list[dict],
                                thresholds: dict,
                                *, logo_action: str | None = None,
                                ) -> dict[str, dict]:
    """L2 跨页重复分析。

    输入：all_texts = [{slide_index, shape_id, text, bbox}]
    返回：shape_id -> {role, preserve_action, detection, is_content_placeholder}
    """
    effective_logo_action = logo_action or _LOGO_DEFAULT_ACTION
    clusters: dict[tuple, list[dict]] = defaultdict(list)
    for item in all_texts:
        t_norm = _normalize_text(item["text"])
        if not t_norm:
            continue
        key = (t_norm, round(item["bbox"]["left"]), round(item["bbox"]["top"]))
        clusters[key].append(item)

    result: dict[str, dict] = {}
    min_repeat = thresholds["cross_slide_repeat_min"]
    for _key, items in clusters.items():
        slides_set = {i["slide_index"] for i in items}
        if len(slides_set) < min_repeat:
            continue
        sample_bbox = items[0]["bbox"]
        top_pct = sample_bbox["top"]
        sample_text = items[0]["text"]
        tlen = len(sample_text)
        if top_pct < 15 and tlen <= 40:
            role = "running_header"
        elif top_pct > 80 and tlen <= 40:
            role = "running_footer"
        elif tlen <= 20:
            role = "logo_text"
        else:
            role = "template_sample_text"

        is_placeholder = _is_content_placeholder(sample_text)
        if role == "logo_text" and sample_text.strip() in _GENERIC_PLACEHOLDER_KEYWORDS:
            is_placeholder = True

        detection = {
            "layer": "L2",
            "signal": f"cross_slide_repeat({len(slides_set)} slides)",
            "confidence": 0.9 if not is_placeholder else 0.5,
            "repeat_slide_count": len(slides_set),
        }
        if role == "template_sample_text":
            preserve_action = "clear_to_empty"
        elif role == "logo_text":
            low = sample_text.strip().lower()
            if any(kw in low for kw in _LOGO_TEXT_KEYWORDS):
                preserve_action = effective_logo_action
            else:
                preserve_action = "keep_original"
        else:
            preserve_action = "keep_original"
        for it in items:
            result[it["shape_id"]] = {
                "role": role,
                "preserve_action": preserve_action,
                "detection": detection,
                "is_content_placeholder": is_placeholder,
            }
    return result


# === story_role 推断（页级叙事角色，保留旧策略并适配）===

def _slide_story_role(idx: int, total: int,
                      content_shape_count: int,
                      sample_texts: list[str]) -> str:
    """根据页位置 + 内容数量 + 文本样本推断 story_role。"""
    if idx == 1:
        return "cover"
    if idx == total:
        return "closing"
    joined = "\n".join(sample_texts)[:500]
    toc_hits = sum(1 for kw in ("目录", "CONTENT", "Contents", "INDEX",
                                "PART", "CHAPTER")
                   if kw.lower() in joined.lower())
    if toc_hits >= 1 and content_shape_count >= 4:
        return "toc"
    if content_shape_count <= 2 and any(_is_chinese(t) for t in sample_texts):
        return "chapter"
    if content_shape_count >= 6:
        return "content_list"
    return "content_detail"


# === content_shape 构建 ===

def _build_content_shape(info: dict, ph_type: str | None,
                         is_placeholder_text: bool,
                         thresholds: dict,
                         *, parent_table_id: str | None = None,
                         table_cell_index: tuple | None = None,
                         ) -> dict:
    """从 shape 信息构建 content_shape 字段（零估算，char_limit 严格等于原文字数）。"""
    bbox = info["bbox"]
    text = info.get("current_text", "") or ""
    char_count = _count_chars(text)
    paragraph_count = max(1, info.get("paragraph_count", 1))
    font_size = info.get("font_size_pt") or 18.0
    body_insets = info.get("body_insets")
    # bbox 估算仅作 debug 字段保留
    estimated_debug = _estimate_char_capacity(bbox, font_size, paragraph_count,
                                              thresholds, body_insets)

    per_para_chars = [_count_chars(p) for p in info.get("paragraphs", [])]
    per_para_sizes = info.get("per_paragraph_font_size_pt", []) or []

    # 逐段 char_limit / font_size（零估算：严格锚定原文字数）
    per_paragraph: list[dict] = []
    if paragraph_count > 1 and per_para_sizes:
        per_paragraph = _compute_per_paragraph_limits(per_para_sizes,
                                                     per_para_chars,
                                                     thresholds)

    # 整 shape char_limit：严格等于 char_count（不乘任何系数）
    # 零字段兜底 1，避免 apply 层被 0 截断
    soft_max = max(1, char_count)
    soft_min = max(1, char_count - 2) if char_count >= 3 else max(1, char_count)
    hard_ceiling = max(1, int(soft_max * 1.15))

    # 若多段总和与整 shape 不一致，以段和为准（更精确）
    if per_paragraph:
        per_sum_soft = sum(p["char_limit"]["max"] for p in per_paragraph)
        per_sum_hard = sum(p["hard_ceiling_chars"] for p in per_paragraph)
        if per_sum_soft > 0:
            soft_max = per_sum_soft
            soft_min = max(1, sum(p["char_limit"]["min"] for p in per_paragraph))
            hard_ceiling = max(hard_ceiling, per_sum_hard)

    return {
        "shape_id": info["shape_id"],
        "debug_name": info.get("debug_name") or info.get("shape_id"),
        "ph_type": ph_type,
        "bbox": bbox,
        "bbox_region": _bbox_region(bbox),
        "original_text": text,
        "char_count": char_count,
        "paragraph_count": paragraph_count,
        "per_paragraph_char_count": per_para_chars,
        "per_paragraph_font_size_pt": per_para_sizes,
        "font_size_pt": round(float(font_size), 2),
        "estimated_capacity": estimated_debug,   # debug only
        "char_limit": {"min": soft_min, "max": soft_max},
        "hard_ceiling_chars": hard_ceiling,
        "per_paragraph": per_paragraph,
        "is_placeholder_text": is_placeholder_text,
        "has_bullet": bool(info.get("has_bullet")),
        "parent_table_id": parent_table_id,
        "table_cell_index": list(table_cell_index) if table_cell_index else None,
    }


# === shape_group 弱提示 ===

def _char_bucket(n: int, edges: list[int]) -> str:
    """按字符数分桶；返回如 '0-10' / '11-30' / '31-80' / '81+'。"""
    sorted_edges = sorted(edges)
    prev = 0
    for e in sorted_edges:
        if n <= e:
            return f"{prev}-{e}" if prev > 0 else f"0-{e}"
        prev = e + 1
    return f"{prev}+"


def _cluster_1d(values: list[float], threshold: float) -> list[list[int]]:
    """1D 贪心聚类：相邻差 < threshold 归同簇。返回索引簇列表。"""
    if not values:
        return []
    pairs = sorted(enumerate(values), key=lambda x: x[1])
    clusters: list[list[int]] = [[pairs[0][0]]]
    last_v = pairs[0][1]
    for idx, v in pairs[1:]:
        if abs(v - last_v) <= threshold:
            clusters[-1].append(idx)
        else:
            clusters.append([idx])
        last_v = v
    return clusters


def _detect_shape_groups(content_shapes: list[dict],
                         thresholds: dict) -> list[dict]:
    """识别 shape_group 弱提示：字数桶 + 桶内 1D 行/列聚类 + 数量 ≥ N。

    仅对 LLM 暴露并列结构提示，不影响 apply 写入。
    """
    min_members = thresholds["shape_group_min_members"]
    var_max = thresholds["shape_group_alignment_var_max_pct"]
    edges = thresholds["char_bucket_edges"]

    buckets: dict[str, list[dict]] = defaultdict(list)
    for sh in content_shapes:
        bucket = _char_bucket(sh["char_count"], edges)
        buckets[bucket].append(sh)

    groups: list[dict] = []
    g_idx = 0
    used_in_group: set[str] = set()

    for bucket, members in buckets.items():
        if len(members) < min_members:
            continue

        # 桶内 top 1D 聚类 → 横向行
        tops = [m["bbox"]["top"] for m in members]
        row_clusters = _cluster_1d(tops, var_max)
        for row_idx_list in row_clusters:
            row_members = [members[i] for i in row_idx_list]
            row_members = [m for m in row_members
                           if m["shape_id"] not in used_in_group]
            if len(row_members) < min_members:
                continue
            row_members.sort(key=lambda m: m["bbox"]["left"])
            avg_chars = sum(m["char_count"] for m in row_members) // len(row_members)
            n = len(row_members)
            hint = (f"{n} 个 ~{avg_chars} 字 shape 横向并列，"
                    f"必须句式一致、字数接近")
            groups.append({
                "group_id": f"g_{g_idx}",
                "member_shape_ids": [m["shape_id"] for m in row_members],
                "char_bucket": bucket,
                "alignment_axis": "horizontal",
                "avg_char_count": avg_chars,
                "member_count": n,
                "group_hint": hint,
            })
            g_idx += 1
            for m in row_members:
                used_in_group.add(m["shape_id"])

        # 桶内剩余 shape 再做 left 1D 聚类 → 纵向列
        remaining = [m for m in members if m["shape_id"] not in used_in_group]
        if len(remaining) < min_members:
            continue
        lefts = [m["bbox"]["left"] for m in remaining]
        col_clusters = _cluster_1d(lefts, var_max)
        for col_idx_list in col_clusters:
            col_members = [remaining[i] for i in col_idx_list]
            if len(col_members) < min_members:
                continue
            col_members.sort(key=lambda m: m["bbox"]["top"])
            avg_chars = sum(m["char_count"] for m in col_members) // len(col_members)
            n = len(col_members)
            hint = (f"{n} 个 ~{avg_chars} 字 shape 纵向并列，"
                    f"必须句式一致、字数接近")
            groups.append({
                "group_id": f"g_{g_idx}",
                "member_shape_ids": [m["shape_id"] for m in col_members],
                "char_bucket": bucket,
                "alignment_axis": "vertical",
                "avg_char_count": avg_chars,
                "member_count": n,
                "group_hint": hint,
            })
            g_idx += 1
            for m in col_members:
                used_in_group.add(m["shape_id"])

    return groups


# === 单 shape 信息提取 ===

def _process_single_shape(shape, ph_type: str | None,
                          slide_idx: int, slide_w_emu: int, slide_h_emu: int,
                          ) -> dict:
    """提取单个 shape 的基础信息（无分类判断）。"""
    bbox = _bbox_pct(shape, slide_w_emu, slide_h_emu)
    sid = _resolve_shape_id(shape, slide_idx)
    info: dict = {
        "shape_id": sid,
        "debug_name": getattr(shape, "name", None) or sid,
        "ph_type": ph_type,
        "bbox": bbox,
        "current_text": "",
        "paragraphs": [],
        "paragraph_count": 0,
        "per_paragraph_font_size_pt": [],
        "font_size_pt": 0.0,
        "has_text_frame": False,
        "has_bullet": False,
        "shape_type": str(shape.shape_type) if shape.shape_type else None,
    }
    if shape.has_text_frame:
        full, paras, weighted, per_para_sizes = _gather_text(shape.text_frame)
        info["current_text"] = full
        info["paragraphs"] = paras
        info["paragraph_count"] = len(paras)
        info["font_size_pt"] = weighted
        info["per_paragraph_font_size_pt"] = [round(s, 2) for s in per_para_sizes]
        info["has_text_frame"] = True
        info["has_bullet"] = _text_frame_has_any_bullet(shape.text_frame)
        info["body_insets"] = _read_body_inset(shape.text_frame)
    return info


# === table cell 拆分（每个 cell 是独立 content_shape） ===

def _iter_table_cells(shape, slide_idx: int,
                      slide_w_emu: int, slide_h_emu: int) -> list[dict]:
    """把 table 的每个 cell 拆为独立 content_shape 候选。"""
    out: list[dict] = []
    try:
        if not shape.has_table:
            return out
    except Exception:
        return out
    parent_id = _resolve_shape_id(shape, slide_idx)
    parent_bbox = _bbox_pct(shape, slide_w_emu, slide_h_emu)
    table = shape.table
    rows = list(table.rows)
    # cell 的精确 bbox 难以获取，采用平均估算
    n_rows = len(rows)
    n_cols = len(rows[0].cells) if n_rows else 0
    if n_rows == 0 or n_cols == 0:
        return out

    cell_w = parent_bbox["w"] / n_cols
    cell_h = parent_bbox["h"] / n_rows
    cell_w_emu = parent_bbox["w_emu"] // n_cols if n_cols else 0
    cell_h_emu = parent_bbox["h_emu"] // n_rows if n_rows else 0

    for r_idx, row in enumerate(rows):
        for c_idx, cell in enumerate(row.cells):
            tf = cell.text_frame
            full, paras, weighted, per_para_sizes = _gather_text(tf)
            cell_body_insets = _read_body_inset(tf)
            cell_bbox = {
                "left": round(parent_bbox["left"] + cell_w * c_idx, 2),
                "top": round(parent_bbox["top"] + cell_h * r_idx, 2),
                "w": round(cell_w, 2),
                "h": round(cell_h, 2),
                "area_pct": round(cell_w * cell_h, 4),
                "left_emu": parent_bbox["left_emu"] + cell_w_emu * c_idx,
                "top_emu": parent_bbox["top_emu"] + cell_h_emu * r_idx,
                "w_emu": cell_w_emu,
                "h_emu": cell_h_emu,
            }
            sid = f"{parent_id}::cell_{r_idx}_{c_idx}"
            out.append({
                "shape_id": sid,
                "debug_name": f"{shape.name or 'Table'}::r{r_idx}c{c_idx}",
                "ph_type": None,
                "bbox": cell_bbox,
                "current_text": full,
                "paragraphs": paras,
                "paragraph_count": len(paras),
                "per_paragraph_font_size_pt": [round(s, 2) for s in per_para_sizes],
                "font_size_pt": weighted,
                "has_text_frame": True,
                "has_bullet": False,
                "shape_type": "TABLE_CELL",
                "body_insets": cell_body_insets,
                "_parent_table_id": parent_id,
                "_table_cell_index": (r_idx, c_idx),
            })
    return out


# === 主解析 ===

def parse(pptx_path: Path, *, thresholds: dict,
          logo_action: str | None = None) -> dict:
    prs = Presentation(str(pptx_path))
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # 第一遍：采集所有 shape 的原始信息
    raw_slides: list[dict] = []
    all_text_items: list[dict] = []  # 供 L2 跨页分析

    for i, slide in enumerate(prs.slides, start=1):
        shapes_info: list[dict] = []
        for shape, _gp in _walk_shapes(slide):
            ph_type = None
            try:
                if shape.is_placeholder and shape.placeholder_format is not None:
                    pft = shape.placeholder_format.type
                    ph_type = pft.name if pft is not None else None
            except Exception:
                pass

            # 表格特殊：拆 cell + 父 shape 也保留以便整体识别
            try:
                if shape.has_table:
                    cells = _iter_table_cells(shape, i, slide_w, slide_h)
                    for cinfo in cells:
                        cinfo["_raw_shape"] = None  # cell 不映射回原 shape
                        shapes_info.append(cinfo)
                        if cinfo["has_text_frame"] and cinfo["current_text"].strip():
                            all_text_items.append({
                                "slide_index": i,
                                "shape_id": cinfo["shape_id"],
                                "text": cinfo["current_text"],
                                "bbox": cinfo["bbox"],
                            })
                    # 父表格本身不再重复添加（其 bbox 已被 cells 覆盖）
                    continue
            except Exception:
                pass

            info = _process_single_shape(shape, ph_type, i, slide_w, slide_h)
            info["_raw_shape"] = shape
            shapes_info.append(info)
            if info["has_text_frame"] and info["current_text"].strip():
                all_text_items.append({
                    "slide_index": i,
                    "shape_id": info["shape_id"],
                    "text": info["current_text"],
                    "bbox": info["bbox"],
                })
        raw_slides.append({"index": i, "shapes": shapes_info})

    # L2 跨页聚类
    l2_results = _cross_slide_repeat_analyze(all_text_items, thresholds,
                                             logo_action=logo_action)

    # 第二遍：分类
    slides_out: list[dict] = []
    vision_ambiguous: list[dict] = []
    vision_pages: set[int] = set()

    for raw in raw_slides:
        i = raw["index"]
        shapes_info = raw["shapes"]
        non_content: list[dict] = []
        content_candidates: list[dict] = []
        sample_texts: list[str] = []

        for info in shapes_info:
            sid = info["shape_id"]
            bbox = info["bbox"]
            ph_type = info["ph_type"]
            raw_shape = info.get("_raw_shape")

            # 0) 特殊占位符（footer/sldNum/dt）—— 一律 keep_original
            sp = _classify_special_placeholder(ph_type)
            if sp is not None:
                non_content.append({
                    "shape_id": sid, "slide_index": i,
                    "current_text": info.get("current_text", ""),
                    "bbox": bbox,
                    **sp,
                })
                continue

            # 1) 图形容器（chart/smartart/picture）—— table 已在第一遍拆 cell 处理
            if raw_shape is not None:
                graphic = _classify_graphic_shape(raw_shape, ph_type, bbox, thresholds)
                if graphic is not None:
                    non_content.append({
                        "shape_id": sid, "slide_index": i,
                        "current_text": info.get("current_text", ""),
                        "bbox": bbox,
                        **graphic,
                    })
                    continue

            # 2) 非文本无特征 shape
            if not info["has_text_frame"]:
                non_content.append({
                    "shape_id": sid, "slide_index": i,
                    "current_text": "",
                    "bbox": bbox,
                    "role": "decoration_shape",
                    "preserve_action": "keep_original",
                    "detection": {"layer": "L1", "signal": "no_text_frame",
                                  "confidence": 1.0},
                })
                continue

            text = info["current_text"]
            font_size = info["font_size_pt"]

            # 3) 空文本 shape：若是占位符（如空 TITLE 框）则纳入 content_shape；否则装饰
            if not text.strip():
                if ph_type in ("TITLE", "CENTER_TITLE", "SUBTITLE", "BODY"):
                    info["_is_empty_placeholder"] = True
                    content_candidates.append(info)
                    continue
                non_content.append({
                    "shape_id": sid, "slide_index": i,
                    "current_text": "", "bbox": bbox,
                    "role": "empty_text_frame",
                    "preserve_action": "keep_original",
                    "detection": {"layer": "L1", "signal": "empty_text",
                                  "confidence": 1.0},
                })
                continue

            # 4) L2 跨页重复
            if sid in l2_results:
                l2 = l2_results[sid]
                if l2.get("is_content_placeholder"):
                    info["_l2_placeholder"] = l2
                    # 放行到 content_candidates
                else:
                    non_content.append({
                        "shape_id": sid, "slide_index": i,
                        "current_text": text, "bbox": bbox,
                        **l2,
                    })
                    continue

            # 5) L1 装饰规则
            if raw_shape is not None:
                l1 = _classify_non_content_l1(raw_shape, ph_type, text, bbox,
                                              font_size, thresholds,
                                              logo_action=logo_action)
                if l1 is not None:
                    non_content.append({
                        "shape_id": sid, "slide_index": i,
                        "current_text": text, "bbox": bbox,
                        **l1,
                    })
                    continue

            # 6) L3 视觉可疑：进 vision_ambiguous，但仍作为内容候选
            signals = _is_vision_ambiguous(text, bbox, ph_type, thresholds)
            if signals is not None:
                vision_ambiguous.append({
                    "shape_id": sid,
                    "slide_index": i,
                    "current_text": text,
                    "bbox_pct": {k: bbox[k] for k in ("left", "top", "w", "h")},
                    "signals": signals,
                    "snapshot_ref": None,
                    "status": "pending",
                    "resolved_role": None,
                    "resolved_preserve_action": None,
                })
                vision_pages.add(i)
                info["_vision_pending"] = True

            content_candidates.append(info)
            if text.strip():
                sample_texts.append(text)

        # 同页内重复短文本检测：≥3 个相同短文本（≤20 字）→ 标 placeholder
        within_page_repeats: dict[str, int] = defaultdict(int)
        for info in content_candidates:
            t = (info.get("current_text") or "").strip()
            if t and len(t) <= 20:
                within_page_repeats[t] += 1
        within_page_placeholders = {
            t for t, n in within_page_repeats.items() if n >= 3
        }

        # 构建 content_shapes
        content_shapes: list[dict] = []
        for info in content_candidates:
            t = (info.get("current_text") or "").strip()
            is_ph = (
                bool(info.get("_l2_placeholder"))
                or bool(info.get("_is_empty_placeholder"))
                or t in within_page_placeholders
            )
            cs = _build_content_shape(
                info, info.get("ph_type"), is_ph, thresholds,
                parent_table_id=info.get("_parent_table_id"),
                table_cell_index=info.get("_table_cell_index"),
            )
            if info.get("_vision_pending"):
                cs["is_vision_pending"] = True
            content_shapes.append(cs)

        # 按空间排序（top, left）—— LLM 自然按视觉顺序读
        content_shapes.sort(key=lambda c: (c["bbox"]["top"], c["bbox"]["left"]))

        # 弱分组提示
        shape_groups = _detect_shape_groups(content_shapes, thresholds)

        # 推断页面叙事角色
        story_role = _slide_story_role(i, len(prs.slides),
                                       len(content_shapes), sample_texts)

        slides_out.append({
            "slide_index": i,
            "story_role": story_role,
            "content_shapes": content_shapes,
            "shape_groups": shape_groups,
            "non_content_shapes": non_content,
        })

    # 汇总
    content_total = sum(len(s["content_shapes"]) for s in slides_out)
    group_total = sum(len(s["shape_groups"]) for s in slides_out)
    non_content_total = sum(len(s["non_content_shapes"]) for s in slides_out)

    return {
        "schema_version": SCHEMA_VERSION,
        "meta": {
            "src_pptx": str(pptx_path),
            "slide_count": len(prs.slides),
            "slide_width_emu": slide_w,
            "slide_height_emu": slide_h,
            "detection_thresholds": thresholds,
            "content_shape_total": content_total,
            "shape_group_total": group_total,
            "non_content_total": non_content_total,
            "vision_ambiguous_total": len(vision_ambiguous),
            "vision_pages": sorted(vision_pages),
        },
        "global_style_guide": None,
        "slides": slides_out,
        "vision_ambiguous": vision_ambiguous,
        "manual_review": [],
    }


# === --resume 回灌（适配新结构）===

def resume_with_vision(story: dict, vision_result: dict) -> dict:
    """把 Prompt B 判定结果回灌到 story（适配 v2 结构）。

    vision_result 结构：
    {
      "results": [
        {"shape_id": "...", "slide_index": N,
         "role": "style_tag" | "content_slot" | ...,
         "preserve_action": "keep_original" | "clear_to_empty" | null,
         "confidence": 0.9, "reason": "..."}
      ]
    }
    """
    results_map: dict[tuple, dict] = {}
    for r in vision_result.get("results", []):
        key = (int(r.get("slide_index", 0)), r.get("shape_id", ""))
        results_map[key] = r

    manual_review: list[dict] = story.get("manual_review", [])
    for amb in story.get("vision_ambiguous", []):
        key = (amb["slide_index"], amb["shape_id"])
        if key not in results_map:
            continue
        r = results_map[key]
        amb["status"] = "resolved"
        amb["resolved_role"] = r.get("role")
        amb["resolved_preserve_action"] = r.get("preserve_action")
        amb["confidence"] = r.get("confidence")
        amb["reason"] = r.get("reason")
        amb["snapshot_ref"] = r.get("snapshot_ref")

        slide_idx = amb["slide_index"]
        slide = next((s for s in story["slides"]
                      if s["slide_index"] == slide_idx), None)
        if slide is None:
            continue

        if r.get("role") == "content_slot":
            # 标记对应 content_shape 解除 vision_pending
            for cs in slide["content_shapes"]:
                if cs["shape_id"] == amb["shape_id"]:
                    cs.pop("is_vision_pending", None)
                    break
        else:
            # non_content：把 content_shape 迁移到 non_content_shapes
            for cs in list(slide["content_shapes"]):
                if cs["shape_id"] == amb["shape_id"]:
                    slide["content_shapes"].remove(cs)
                    slide["non_content_shapes"].append({
                        "shape_id": amb["shape_id"],
                        "slide_index": slide_idx,
                        "current_text": cs.get("original_text", ""),
                        "bbox": cs.get("bbox"),
                        "role": r.get("role"),
                        "preserve_action": r.get("preserve_action") or "keep_original",
                        "detection": {
                            "layer": "L3_vision",
                            "signal": "prompt_b_resume",
                            "confidence": r.get("confidence"),
                            "reason": r.get("reason"),
                        },
                    })
                    break

    for amb in story.get("vision_ambiguous", []):
        if amb["status"] == "pending":
            amb["status"] = "failed"
            amb["reason"] = "no_vision_result_provided"

    story["manual_review"] = manual_review
    story["meta"]["non_content_total"] = sum(
        len(s["non_content_shapes"]) for s in story["slides"])
    story["meta"]["content_shape_total"] = sum(
        len(s["content_shapes"]) for s in story["slides"])
    return story


# === debug dump（人类可读清单）===

def render_debug_md(story: dict) -> str:
    """生成 template_story_debug.md 内容。"""
    lines = [
        f"# Template Story Debug (schema v{story['schema_version']})",
        "",
        f"- src: `{story['meta'].get('src_pptx')}`",
        f"- slides: {story['meta'].get('slide_count')}",
        f"- content_shape_total: {story['meta'].get('content_shape_total')}",
        f"- shape_group_total: {story['meta'].get('shape_group_total')}",
        f"- non_content_total: {story['meta'].get('non_content_total')}",
        f"- vision_ambiguous_total: {story['meta'].get('vision_ambiguous_total')}",
        "",
    ]
    for s in story.get("slides", []):
        lines.append(
            f"## Slide {s['slide_index']} - story_role={s['story_role']}")
        lines.append("")
        cs_list = s.get("content_shapes", [])
        if cs_list:
            lines.append("### content_shapes")
            lines.append("")
            lines.append("| shape_id | region | chars | font | cap | limit | "
                         "hard | paras | ph | placeholder | text |")
            lines.append("|---|---|---|---|---|---|---|---|---|---|---|")
            for cs in cs_list:
                txt = (cs.get("original_text") or "").replace("\n", " ").strip()
                if len(txt) > 24:
                    txt = txt[:24] + "..."
                limit = cs.get("char_limit", {})
                lines.append(
                    f"| `{cs['shape_id']}` "
                    f"| {cs.get('bbox_region', '?')} "
                    f"| {cs.get('char_count', 0)} "
                    f"| {cs.get('font_size_pt', 0)} "
                    f"| {cs.get('estimated_capacity', 0)} "
                    f"| {limit.get('min', '?')}-{limit.get('max', '?')} "
                    f"| {cs.get('hard_ceiling_chars', '?')} "
                    f"| {cs.get('paragraph_count', 1)} "
                    f"| {cs.get('ph_type') or ''} "
                    f"| {'Y' if cs.get('is_placeholder_text') else ''} "
                    f"| {txt} |"
                )
            lines.append("")
            # per_paragraph 细节
            multi = [cs for cs in cs_list if cs.get("per_paragraph")]
            if multi:
                lines.append("#### per_paragraph")
                lines.append("")
                lines.append("| shape_id | idx | font | limit | hard | emphasis |")
                lines.append("|---|---|---|---|---|---|")
                for cs in multi:
                    for pp in cs["per_paragraph"]:
                        pl = pp.get("char_limit", {})
                        lines.append(
                            f"| `{cs['shape_id']}` "
                            f"| {pp.get('idx')} "
                            f"| {pp.get('font_size_pt')} "
                            f"| {pl.get('min', '?')}-{pl.get('max', '?')} "
                            f"| {pp.get('hard_ceiling_chars', '?')} "
                            f"| {'Y' if pp.get('is_emphasis') else ''} |"
                        )
                lines.append("")
        sg_list = s.get("shape_groups", [])
        if sg_list:
            lines.append("### shape_groups")
            lines.append("")
            for g in sg_list:
                lines.append(
                    f"- **{g['group_id']}** "
                    f"({g['member_count']} members, "
                    f"axis={g['alignment_axis']}, "
                    f"bucket={g['char_bucket']}, "
                    f"avg_chars={g['avg_char_count']}): "
                    f"{g.get('group_hint', '')}")
                for sid in g["member_shape_ids"]:
                    lines.append(f"  - `{sid}`")
            lines.append("")
        nc_list = s.get("non_content_shapes", [])
        if nc_list:
            lines.append("### non_content_shapes")
            lines.append("")
            lines.append("| shape_id | role | action | text |")
            lines.append("|---|---|---|---|")
            for nc in nc_list:
                txt = (nc.get("current_text") or "").replace("\n", " ").strip()
                if len(txt) > 24:
                    txt = txt[:24] + "..."
                lines.append(
                    f"| `{nc['shape_id']}` "
                    f"| {nc.get('role', '?')} "
                    f"| {nc.get('preserve_action', '?')} "
                    f"| {txt} |"
                )
            lines.append("")
        lines.append("")
    return "\n".join(lines)


# === CLI ===

def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("pptx", type=Path, nargs="?")
    ap.add_argument("--out", required=True, type=Path)
    ap.add_argument("--resume", type=Path, default=None,
                    help="vision_result.json；提供时从 --story 读入已有 story 并回灌")
    ap.add_argument("--story", type=Path, default=None,
                    help="--resume 模式下的已有 template_story.json")
    ap.add_argument("--cross-slide-repeat-min", type=int,
                    default=DEFAULT_THRESHOLDS["cross_slide_repeat_min"])
    ap.add_argument("--logo-action", choices=("keep_original", "clear_to_empty"),
                    default=_LOGO_DEFAULT_ACTION,
                    help="LOGO 占位词命中后的默认行为（默认 keep_original）")
    ap.add_argument("--shape-group-min", type=int,
                    default=DEFAULT_THRESHOLDS["shape_group_min_members"],
                    help="shape_group 弱提示的最小成员数（默认 3）")
    ap.add_argument("--no-debug-dump", action="store_true",
                    help="跳过 <out>.debug.md 生成")
    args = ap.parse_args()

    thresholds = dict(DEFAULT_THRESHOLDS)
    thresholds["cross_slide_repeat_min"] = args.cross_slide_repeat_min
    thresholds["shape_group_min_members"] = args.shape_group_min

    if args.resume is not None:
        story_path = args.story or args.out
        if not story_path.exists():
            print(f"ERROR: --resume 需要 --story 指向已有 template_story.json: "
                  f"{story_path}", file=sys.stderr)
            return 1
        story = json.loads(story_path.read_text(encoding="utf-8"))
        vision_result = json.loads(args.resume.read_text(encoding="utf-8"))
        updated = resume_with_vision(story, vision_result)
        args.out.parent.mkdir(parents=True, exist_ok=True)
        args.out.write_text(json.dumps(updated, ensure_ascii=False, indent=2),
                            encoding="utf-8")
        if not args.no_debug_dump:
            debug_path = args.out.with_suffix(args.out.suffix + ".debug.md")
            debug_path.write_text(render_debug_md(updated), encoding="utf-8")
        resolved = sum(1 for a in updated.get("vision_ambiguous", [])
                       if a["status"] == "resolved")
        print(json.dumps({
            "ok": True, "out": str(args.out),
            "resume": True,
            "vision_resolved": resolved,
            "vision_failed": sum(1 for a in updated.get("vision_ambiguous", [])
                                 if a["status"] == "failed"),
            "manual_review_count": len(updated.get("manual_review", [])),
        }, ensure_ascii=False))
        return 0

    if args.pptx is None or not args.pptx.exists():
        print(f"ERROR: pptx not found: {args.pptx}", file=sys.stderr)
        return 1

    story = parse(args.pptx, thresholds=thresholds,
                  logo_action=args.logo_action)
    args.out.parent.mkdir(parents=True, exist_ok=True)
    args.out.write_text(json.dumps(story, ensure_ascii=False, indent=2),
                        encoding="utf-8")
    if not args.no_debug_dump:
        debug_path = args.out.with_suffix(args.out.suffix + ".debug.md")
        debug_path.write_text(render_debug_md(story), encoding="utf-8")

    meta = story["meta"]
    print(json.dumps({
        "ok": True, "out": str(args.out),
        "slide_count": meta["slide_count"],
        "content_shape_total": meta["content_shape_total"],
        "shape_group_total": meta["shape_group_total"],
        "non_content_total": meta["non_content_total"],
        "vision_ambiguous_total": meta["vision_ambiguous_total"],
        "vision_pages": meta["vision_pages"],
        "schema_version": SCHEMA_VERSION,
    }, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    sys.exit(main())
