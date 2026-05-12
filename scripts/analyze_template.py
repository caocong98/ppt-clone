"""[传统脚本] 解析 PPT 模板：占位区清单 + slide 启发式角色 + 容量约束 + 非文本块标记。

说明：新流程（narrative projection）**不再依赖本脚本**的 fill_priority/heuristic_role。
这些判定全部迁移到 parse_template_story.py（story_role + narrative_slots）。本脚本
现在主要作用是提供给颜色 / 渲染阶段的几何、shape 清单，以及兼容旧的 mode_a 流程。

v3 新增（防过填，现由 parse_template_story 接管）：
- fill_priority ∈ {required, optional, decoration, numeric_marker}
- is_decoration_marker：纯数字编号 / 装饰性 LOGO 占位符识别
- current_text_len、area_pct：供 LLM 与下游 lint 判断
- schema_version=3

v2 增强（保留）：
- 字号继承：run -> paragraph defRPr -> layout 同 idx 占位 -> master 同 idx 占位 -> 默认 18
- group 展开：递归进入 MSO_SHAPE_TYPE.GROUP，扁平化所有内部文本框
- 容量公式可配：--safety / --zh-factor / --en-factor / --line-spacing

CLI:
    python analyze_template.py <pptx> --out template_spec.json
        [--safety 0.9] [--zh-factor 0.95] [--en-factor 0.55] [--line-spacing 1.2]
        [--no-expand-groups]
        [--decoration-area-pct 0.3] [--decoration-font-pt 9]
        [--marker-area-pct 1.0]
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SCHEMA_VERSION = 3

# 装饰性 marker 文本识别
_NUMBER_MARKER_RE = re.compile(r"^\s*0?\d{1,2}\s*[\.、:：\-/]?\s*$")
_PURE_PUNCT_RE = re.compile(r"^[\W_\s]+$", re.UNICODE)
_DECO_KEYWORDS = {
    "your logo", "logo", "icon", "图标", "logo here", "your name",
    "company name", "公司名称",
}

TEXT_PH_KINDS = {"TITLE", "CENTER_TITLE", "SUBTITLE", "BODY", "OBJECT", "VERTICAL_TITLE", "VERTICAL_BODY"}
SKIP_PH_KINDS = {"FOOTER", "HEADER", "SLIDE_NUMBER", "DATE"}
NON_TEXT_PH_KINDS = {"PICTURE", "BITMAP", "MEDIA", "CHART", "TABLE", "ORG_CHART"}


def _emu_to_pt(emu: int) -> float:
    return emu / 12700.0


def _is_chinese(s: str) -> bool:
    return any("\u4e00" <= ch <= "\u9fff" for ch in s)


def _estimate_max_chars(width_pt: float, font_size_pt: float, text_sample: str,
                        zh_factor: float, en_factor: float) -> int:
    if font_size_pt <= 0:
        font_size_pt = 18
    factor = zh_factor if _is_chinese(text_sample) or not text_sample else en_factor
    if not text_sample:
        factor = zh_factor  # 模板大概率是中文场景，保守估计
    return max(1, int(width_pt / (font_size_pt * factor)))


def _estimate_max_lines(height_pt: float, font_size_pt: float, line_spacing: float) -> int:
    if font_size_pt <= 0:
        font_size_pt = 18
    return max(1, int(height_pt / (font_size_pt * line_spacing)))


def _is_decoration_marker(text: str) -> bool:
    """判断文本是否是装饰性编号 / LOGO 占位符。"""
    if not text:
        return False
    t = text.strip()
    if not t:
        return False
    if len(t) <= 4 and _NUMBER_MARKER_RE.match(t):
        return True
    if t.lower() in _DECO_KEYWORDS:
        return True
    if len(t) <= 6 and _PURE_PUNCT_RE.match(t):
        return True
    return False


def _decide_fill_priority(
    *, ph_type: str | None, role: str, text: str, text_len: int,
    area_pct: float, font_size_pt: float, is_marker: bool,
    decoration_area_pct: float, decoration_font_pt: float, marker_area_pct: float,
) -> str:
    """fill_priority 决策：

    - numeric_marker：装饰性编号/LOGO 类（且面积小）
    - decoration：极小面积 / 字号过小 / 极短文本+小面积
    - required：标题占位符、或大面积 + 正常字号
    - 否则 optional
    """
    if is_marker and area_pct < marker_area_pct:
        return "numeric_marker"
    if area_pct < decoration_area_pct:
        return "decoration"
    if font_size_pt > 0 and font_size_pt < decoration_font_pt:
        return "decoration"
    if text_len <= 2 and area_pct < 0.6:
        return "decoration"
    if ph_type in ("TITLE", "CENTER_TITLE", "SUBTITLE"):
        return "required"
    if role in ("title", "subtitle"):
        return "required"
    if area_pct >= 4.0 and font_size_pt >= 14:
        return "required"
    return "optional"


def _detect_role_for_text(font_size_pt: float) -> str:
    if font_size_pt >= 32:
        return "title"
    if font_size_pt >= 22:
        return "subtitle"
    return "body"


def _heuristic_slide_role(idx: int, total: int, placeholders: list[dict], shapes_count: int) -> str:
    titles = [p for p in placeholders if p["role"] == "title"]
    bodies = [p for p in placeholders if p["role"] in ("body", "bullets")]

    if idx == 1:
        return "cover"
    if idx == total:
        return "thanks" if shapes_count <= 3 else "summary"
    if idx == 2 and bodies and any(
        "目录" in (p.get("current_text") or "") or "contents" in (p.get("current_text") or "").lower()
        for p in placeholders
    ):
        return "toc"
    if titles and not bodies and shapes_count <= 5:
        return "section"
    return "content"


def _run_size_pt(run) -> float | None:
    sz = run.font.size
    return sz.pt if sz is not None else None


def _paragraph_def_size_pt(paragraph) -> float | None:
    """从 a:pPr/a:defRPr@sz 读段落默认字号（百分点 *100）。"""
    try:
        pPr = paragraph._pPr
        if pPr is None:
            return None
        defRPr = pPr.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr"
        )
        if defRPr is None:
            return None
        sz = defRPr.get("sz")
        if sz is None:
            return None
        return int(sz) / 100.0
    except Exception:
        return None


def _placeholder_inherited_size(slide_shape, layout, master) -> float | None:
    """对 placeholder shape，按 idx 在 layout/master 上找同槽位的字号。"""
    if not slide_shape.is_placeholder or slide_shape.placeholder_format is None:
        return None
    idx = slide_shape.placeholder_format.idx
    if idx is None:
        return None

    for src in (layout, master):
        if src is None:
            continue
        try:
            for ph in src.placeholders:
                if ph.placeholder_format is not None and ph.placeholder_format.idx == idx:
                    if ph.has_text_frame:
                        for p in ph.text_frame.paragraphs:
                            for r in p.runs:
                                v = _run_size_pt(r)
                                if v:
                                    return v
                            v = _paragraph_def_size_pt(p)
                            if v:
                                return v
        except Exception:
            continue
    return None


def _gather_text(text_frame, fallback_size: float) -> tuple[str, list[str], float]:
    paragraphs = []
    sizes: list[float] = []
    for p in text_frame.paragraphs:
        text = "".join(r.text or "" for r in p.runs)
        if not text and p.text:
            text = p.text
        paragraphs.append(text)
        for r in p.runs:
            sz = _run_size_pt(r)
            if sz:
                sizes.append(sz)
        # 段落 defRPr 兜底
        if not sizes:
            v = _paragraph_def_size_pt(p)
            if v:
                sizes.append(v)
    full_text = "\n".join(paragraphs)
    dominant_size = max(sizes) if sizes else fallback_size
    return full_text, paragraphs, dominant_size


def _walk_shapes(container, expand_groups: bool):
    """yield 所有 shape，可选展开 group。"""
    for shape in container.shapes if hasattr(container, "shapes") else container:
        if expand_groups and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape, expand_groups=True)
        else:
            yield shape


def _process_shape(shape, slide, layout, master, slide_w_emu, slide_h_emu,
                   safety: float, zh_factor: float, en_factor: float, line_spacing: float,
                   decoration_area_pct: float, decoration_font_pt: float, marker_area_pct: float,
                   group_path: str = "") -> tuple[dict | None, dict | None]:
    """返回 (placeholder_dict | None, non_text_dict | None)。"""
    shape_id = shape.name if not group_path else f"{group_path}/{shape.name}"

    ph_type = None
    if shape.is_placeholder and shape.placeholder_format is not None:
        pft = shape.placeholder_format.type
        ph_type = pft.name if pft is not None else None

    try:
        left = shape.left or 0
        top = shape.top or 0
        width = shape.width or 0
        height = shape.height or 0
    except Exception:
        left = top = width = height = 0

    if ph_type in SKIP_PH_KINDS:
        return None, None

    if ph_type in NON_TEXT_PH_KINDS or shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.MEDIA):
        return None, {
            "shape_id": shape_id,
            "kind": "picture" if (ph_type in ("PICTURE", "BITMAP")
                                  or shape.shape_type == MSO_SHAPE_TYPE.PICTURE) else (ph_type or "media").lower(),
            "bbox": {"x": left, "y": top, "w": width, "h": height},
            "action": "skip",
        }
    if shape.has_chart or shape.shape_type == MSO_SHAPE_TYPE.CHART:
        return None, {
            "shape_id": shape_id, "kind": "chart",
            "bbox": {"x": left, "y": top, "w": width, "h": height},
            "action": "theme_only_no_text",
        }
    if shape.has_table:
        return None, {
            "shape_id": shape_id, "kind": "table",
            "bbox": {"x": left, "y": top, "w": width, "h": height},
            "action": "theme_only_no_text",
        }

    if not shape.has_text_frame:
        return None, None

    inherited = _placeholder_inherited_size(shape, layout, master)
    fallback = inherited if inherited else 18.0
    full_text, paragraphs, dom_size = _gather_text(shape.text_frame, fallback)

    width_pt = _emu_to_pt(width)
    height_pt = _emu_to_pt(height)

    current_bullets = [p for p in paragraphs if p.strip()]
    role = _detect_role_for_text(dom_size)
    if ph_type in ("TITLE", "CENTER_TITLE"):
        role = "title"
    elif ph_type == "SUBTITLE":
        role = "subtitle"
    elif ph_type in ("BODY", "OBJECT"):
        role = "bullets" if len(current_bullets) > 1 else "body"

    raw_max_chars = _estimate_max_chars(width_pt, dom_size, full_text, zh_factor, en_factor)
    raw_max_lines = _estimate_max_lines(height_pt, dom_size, line_spacing)
    safe_max_chars = max(1, int(raw_max_chars * safety))
    safe_max_lines = max(1, int(raw_max_lines * safety))
    max_bullets = max(1, len(current_bullets)) if role == "bullets" else 1
    max_chars_per_bullet = safe_max_chars

    current_text_len = len(full_text.strip())
    slide_area = max(1, slide_w_emu * slide_h_emu)
    area_pct = (width * height) / slide_area * 100.0
    is_marker = _is_decoration_marker(full_text)
    fill_priority = _decide_fill_priority(
        ph_type=ph_type, role=role, text=full_text, text_len=current_text_len,
        area_pct=area_pct, font_size_pt=dom_size, is_marker=is_marker,
        decoration_area_pct=decoration_area_pct, decoration_font_pt=decoration_font_pt,
        marker_area_pct=marker_area_pct,
    )
    # 期望长度区间：以原文长度为锚，避免大幅偏移
    if current_text_len > 0:
        target_min = max(1, int(current_text_len * 0.5))
        target_max = max(target_min + 1, int(current_text_len * 1.6))
    else:
        target_min, target_max = 0, max(1, safe_max_chars)
    target_max = min(target_max, safe_max_chars)

    return {
        "shape_id": shape_id,
        "ph_type": ph_type,
        "role": role,
        "current_text": full_text,
        "current_text_len": current_text_len,
        "current_bullets": current_bullets,
        "paragraph_count": len(paragraphs),
        "shape_width_emu": width,
        "shape_height_emu": height,
        "shape_left_emu": left,
        "shape_top_emu": top,
        "area_pct": round(area_pct, 4),
        "font_size_pt": dom_size,
        "font_size_source": "run" if dom_size != fallback else ("inherited" if inherited else "fallback"),
        "raw_max_chars": raw_max_chars,
        "raw_max_lines": raw_max_lines,
        "max_chars": safe_max_chars,
        "max_lines": safe_max_lines,
        "max_bullets": max_bullets,
        "max_chars_per_bullet": max_chars_per_bullet,
        "target_chars_min": target_min,
        "target_chars_max": target_max,
        "is_decoration_marker": is_marker,
        "fill_priority": fill_priority,
        "in_group": bool(group_path),
        "group_path": group_path or None,
        "fillable": fill_priority not in ("numeric_marker",),
    }, None


def _walk_with_path(container, expand_groups: bool, path: str = ""):
    for shape in container.shapes if hasattr(container, "shapes") else container:
        if expand_groups and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            sub_path = f"{path}/{shape.name}" if path else shape.name
            yield from _walk_with_path(shape, expand_groups=True, path=sub_path)
        else:
            yield shape, path


def analyze(pptx_path: Path, *, safety: float, zh_factor: float, en_factor: float,
            line_spacing: float, expand_groups: bool,
            decoration_area_pct: float = 0.3, decoration_font_pt: float = 9.0,
            marker_area_pct: float = 1.0) -> dict:
    prs = Presentation(str(pptx_path))
    slide_w_emu = prs.slide_width
    slide_h_emu = prs.slide_height

    slides_data = []
    for i, slide in enumerate(prs.slides, start=1):
        layout = slide.slide_layout if hasattr(slide, "slide_layout") else None
        master = layout.slide_master if (layout is not None and hasattr(layout, "slide_master")) else None
        layout_name = layout.name if layout else ""

        placeholders = []
        non_text_blocks = []
        groups_seen: list[dict] = []

        # 顶层 shape 同时也单独标 group bbox（便于 mask 等下游用）
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    groups_seen.append({
                        "shape_id": shape.name,
                        "kind": "group",
                        "bbox": {"x": shape.left or 0, "y": shape.top or 0,
                                 "w": shape.width or 0, "h": shape.height or 0},
                        "action": "expanded" if expand_groups else "skip",
                    })
                except Exception:
                    pass

        for shape, group_path in _walk_with_path(slide, expand_groups=expand_groups):
            ph, nt = _process_shape(
                shape, slide, layout, master, slide_w_emu, slide_h_emu,
                safety=safety, zh_factor=zh_factor, en_factor=en_factor,
                line_spacing=line_spacing,
                decoration_area_pct=decoration_area_pct,
                decoration_font_pt=decoration_font_pt,
                marker_area_pct=marker_area_pct,
                group_path=group_path,
            )
            if ph is not None:
                placeholders.append(ph)
            if nt is not None:
                non_text_blocks.append(nt)

        non_text_blocks.extend(groups_seen)

        heuristic_role = _heuristic_slide_role(i, len(prs.slides), placeholders, len(slide.shapes))

        slides_data.append({
            "index": i,
            "layout_name": layout_name,
            "heuristic_role": heuristic_role,
            "shape_count": len(slide.shapes),
            "placeholders": placeholders,
            "non_text_blocks": non_text_blocks,
        })

    # 顶层统计：每页 fill_priority 分布、decoration/marker 数量
    summary = {"required": 0, "optional": 0, "decoration": 0, "numeric_marker": 0}
    for s in slides_data:
        for ph in s["placeholders"]:
            fp = ph.get("fill_priority", "optional")
            summary[fp] = summary.get(fp, 0) + 1

    return {
        "schema_version": SCHEMA_VERSION,
        "slide_width_emu": slide_w_emu,
        "slide_height_emu": slide_h_emu,
        "slide_count": len(prs.slides),
        "config": {
            "safety": safety, "zh_factor": zh_factor, "en_factor": en_factor,
            "line_spacing": line_spacing, "expand_groups": expand_groups,
            "decoration_area_pct": decoration_area_pct,
            "decoration_font_pt": decoration_font_pt,
            "marker_area_pct": marker_area_pct,
        },
        "fill_priority_summary": summary,
        "slides": slides_data,
    }


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("pptx", type=Path)
    ap.add_argument("--out", required=True, type=Path)
    ap.add_argument("--safety", type=float, default=0.9)
    ap.add_argument("--zh-factor", type=float, default=0.95)
    ap.add_argument("--en-factor", type=float, default=0.55)
    ap.add_argument("--line-spacing", type=float, default=1.2)
    ap.add_argument("--no-expand-groups", action="store_true")
    ap.add_argument("--decoration-area-pct", type=float, default=0.3,
                    help="<= 该百分比面积 -> decoration")
    ap.add_argument("--decoration-font-pt", type=float, default=9.0,
                    help="< 该字号 -> decoration")
    ap.add_argument("--marker-area-pct", type=float, default=1.0,
                    help="装饰编号/LOGO 文本且面积 < 该百分比 -> numeric_marker")
    args = ap.parse_args()

    if not args.pptx.exists():
        print(f"ERROR: pptx not found: {args.pptx}", file=sys.stderr)
        return 1

    result = analyze(
        args.pptx,
        safety=args.safety, zh_factor=args.zh_factor, en_factor=args.en_factor,
        line_spacing=args.line_spacing, expand_groups=not args.no_expand_groups,
        decoration_area_pct=args.decoration_area_pct,
        decoration_font_pt=args.decoration_font_pt,
        marker_area_pct=args.marker_area_pct,
    )
    args.out.parent.mkdir(parents=True, exist_ok=True)
    args.out.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps({
        "ok": True, "out": str(args.out),
        "slide_count": result["slide_count"],
        "placeholder_total": sum(len(s["placeholders"]) for s in result["slides"]),
        "fill_priority_summary": result["fill_priority_summary"],
        "schema_version": SCHEMA_VERSION,
        "expand_groups": not args.no_expand_groups,
    }, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    sys.exit(main())
